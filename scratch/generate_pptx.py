# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = "c:\\Users\\alexc\\OneDrive\\文件\\WikiLLM\\WikiLLM\\outputs"

# --- Shared Colors ---
C_WHITE = RGBColor(255, 255, 255)
C_DARK_BG = RGBColor(10, 25, 47)
C_DARK_TEXT = RGBColor(40, 40, 40)
C_LLM_BG = RGBColor(25, 40, 70)       # Darker blue for base
C_BCR = RGBColor(71, 85, 105)         # Slate
C_CRM = RGBColor(51, 65, 85)          # Darker Slate
C_CLM = RGBColor(74, 222, 128)        # Neon Green
C_BPM = RGBColor(56, 189, 248)        # Neon Cyan
C_ESIGN = RGBColor(251, 146, 60)      # Neon Orange
C_KM = RGBColor(244, 114, 182)        # Neon Pink

# Helper to draw the native 16:9 vector framework
def draw_native_framework(slide, left_start, top_start, total_width, total_height, is_dark_theme):
    # Base Box: Local LLM (bottom)
    llm_height = Inches(0.8)
    llm_top = top_start + total_height - llm_height
    
    llm_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_start, llm_top, total_width, llm_height)
    llm_rect.fill.solid()
    llm_rect.fill.fore_color.rgb = C_LLM_BG if is_dark_theme else RGBColor(220, 230, 245)
    llm_rect.line.color.rgb = C_BPM
    
    tf = llm_rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Local LLM AI Model (Central Enterprise Brain)"
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_WHITE if is_dark_theme else C_DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # The 6 Pillars
    pillar_gap = Inches(0.15)
    pillar_width = (total_width - (5 * pillar_gap)) / 6
    pillar_height = total_height - llm_height - Inches(0.3) # leaves gap above LLM
    
    pillars_data = [
        ("BCR", "OCR Sensing", C_BCR, ["Business Card OCR", "Auto Data Sync"]),
        ("CRM", "Pipedrive", C_CRM, ["Contacts & Leads", "Deal Pipeline"]),
        ("CLM", "doc-cowork", C_CLM, ["MS-Word / Google Doc", "Collaborative Editing"]),
        ("BPM", "Workflow", C_BPM, ["All Departments", "AI-Review Routing"]),
        ("ESign", "BreezySign", C_ESIGN, ["Sales & Admin Sign", "Call BZS API"]),
        ("KM", "Files Manager", C_KM, ["RAW Repository", "Obsidian Graphify"])
    ]
    
    current_left = left_start
    for title, subtitle, color, items in pillars_data:
        # Pillar container
        prect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, current_left, top_start, pillar_width, pillar_height)
        prect.fill.solid()
        if is_dark_theme:
            # Dark theme uses hollow pillars with colored borders
            prect.fill.background()
            prect.line.color.rgb = color
            prect.line.width = Pt(1.5)
            text_color = color
        else:
            # Light theme uses solid colored pillars
            prect.fill.fore_color.rgb = color
            prect.line.fill.background()
            text_color = C_WHITE if title != "CLM" and title != "BPM" else C_DARK_TEXT
            
        # Pillar Title Box
        tbox = slide.shapes.add_textbox(current_left, top_start + Inches(0.1), pillar_width, Inches(0.8))
        pt = tbox.text_frame.paragraphs[0]
        pt.text = title
        pt.font.name = 'Arial'
        pt.font.size = Pt(22)
        pt.font.bold = True
        pt.font.color.rgb = text_color
        pt.alignment = PP_ALIGN.CENTER
        
        pst = tbox.text_frame.add_paragraph()
        pst.text = subtitle
        pst.font.name = 'Arial'
        pst.font.size = Pt(12)
        pst.font.color.rgb = text_color
        pst.alignment = PP_ALIGN.CENTER
        
        # Sub-blocks within pillar
        block_h = Inches(0.8)
        block_gap = Inches(0.15)
        b_top = top_start + Inches(1.2)
        for idx, item in enumerate(items):
            b_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, current_left + Inches(0.1), b_top, pillar_width - Inches(0.2), block_h)
            b_rect.fill.solid()
            b_rect.fill.fore_color.rgb = color
            b_rect.line.fill.background()
            
            pb = b_rect.text_frame.paragraphs[0]
            pb.text = item
            pb.font.name = 'Arial'
            pb.font.size = Pt(13)
            pb.font.bold = True
            pb.font.color.rgb = C_DARK_BG if is_dark_theme else C_WHITE
            pb.alignment = PP_ALIGN.CENTER
            
            b_top += block_h + block_gap
            
        # Connecting arrow to LLM base
        conn = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, current_left + (pillar_width/2) - Inches(0.1), top_start + pillar_height + Inches(0.05), Inches(0.2), Inches(0.2))
        conn.fill.solid()
        conn.fill.fore_color.rgb = color
        conn.line.fill.background()
        
        current_left += pillar_width + pillar_gap


# ---------------------------------------------------------
# 1. GENERATE GENERAL NEON DARK-THEME PRESENTATION
# ---------------------------------------------------------
def create_general_neon_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    COLOR_BG = RGBColor(10, 25, 47)
    blank_layout = prs.slide_layouts[6]
    
    def set_dark_bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLOR_BG
        
    def add_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.13), Inches(0.8))
        p = title_box.text_frame.paragraphs[0]
        p.text = title_text
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C_BPM
        return title_box

    # --- Cover ---
    slide1 = prs.slides.add_slide(blank_layout)
    set_dark_bg(slide1)
    cb = slide1.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12.13), Inches(4.5))
    tf1 = cb.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "BreezyBrain (好好腦)"
    p1.font.size, p1.font.bold, p1.font.color.rgb = Pt(54), True, C_ESIGN
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "下一代 AI 企業工作流操作系統 (16:9 向量架構版)"
    p1_sub.font.size, p1_sub.font.bold, p1_sub.font.color.rgb = Pt(28), True, C_WHITE
    p1_sub.space_before = Pt(10)

    # --- FRAMEWORK DIAGRAM (NATIVE) ---
    slide_fw = prs.slides.add_slide(blank_layout)
    set_dark_bg(slide_fw)
    add_header(slide_fw, "一、 BreezyBrain 六大支柱系統架構圖 (Native 16:9 Vector)")
    
    # Draw native vector framework in the center area
    draw_native_framework(slide_fw, Inches(1.0), Inches(1.5), Inches(11.33), Inches(5.3), True)

    prs.save(os.path.join(OUTPUT_DIR, "BreezyBrain_General_Edition_v3.pptx"))
    print("Generated: BreezyBrain_General_Edition_v3.pptx")


# ---------------------------------------------------------
# 2. GENERATE PENPOWER CIS PRESENTATION
# ---------------------------------------------------------
def create_penpower_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    COLOR_PRI_BLUE = RGBColor(0, 156, 223)
    COLOR_LBG = RGBColor(227, 242, 250)
    blank = prs.slide_layouts[6]
    
    def set_bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color
        
    def add_rect(slide, l, t, w, h, color):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        
    def add_out_rect(slide, l, t, w, h, color):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        rect.fill.background()
        rect.line.color.rgb = color
        rect.line.width = Pt(1.5)

    # --- Cover ---
    slide1 = prs.slides.add_slide(blank)
    set_bg(slide1, COLOR_LBG)
    btop, bh = Inches(1.8), Inches(4.2)
    add_rect(slide1, 0, btop, Inches(13.33), bh, COLOR_PRI_BLUE)
    
    # Text
    cb = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(10), Inches(3.0))
    p1 = cb.text_frame.paragraphs[0]
    p1.text = "BreezyBrain (好好腦)"
    p1.font.size, p1.font.bold, p1.font.color.rgb = Pt(54), True, C_WHITE
    p1_sub = cb.text_frame.add_paragraph()
    p1_sub.text = "下一代 AI 企業工作流操作系統 (16:9 向量架構版)"
    p1_sub.font.size, p1_sub.font.color.rgb, p1_sub.space_before = Pt(24), C_WHITE, Pt(15)

    def add_content_slide(title_text):
        slide = prs.slides.add_slide(blank)
        set_bg(slide, C_WHITE)
        m = Inches(0.4)
        add_out_rect(slide, m, m, Inches(13.33) - 2*m, Inches(7.5) - 2*m, COLOR_PRI_BLUE)
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(0.8))
        tb.text_frame.paragraphs[0].text = title_text
        tb.text_frame.paragraphs[0].font.size, tb.text_frame.paragraphs[0].font.color.rgb = Pt(28), COLOR_PRI_BLUE
        tb.text_frame.paragraphs[0].font.bold = True
        return slide

    # --- FRAMEWORK DIAGRAM (NATIVE) ---
    s_fw = add_content_slide("一、 BreezyBrain 六大支柱系統架構圖 (Native 16:9 Vector)")
    
    # Draw native vector framework in the center area
    draw_native_framework(s_fw, Inches(1.0), Inches(1.6), Inches(11.33), Inches(5.0), False)

    prs.save(os.path.join(OUTPUT_DIR, "BreezyBrain_PenPower_Edition_v3.pptx"))
    print("Generated: BreezyBrain_PenPower_Edition_v3.pptx")

if __name__ == "__main__":
    print("Starting native 16:9 vector framework PPTX generation...")
    create_general_neon_presentation()
    create_penpower_presentation()
    print("All native presentations generated successfully.")
