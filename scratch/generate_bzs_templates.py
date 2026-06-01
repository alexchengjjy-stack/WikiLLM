# -*- coding: utf-8 -*-
import os
import subprocess
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. 定義路徑
HTML_FILE = r"c:\Users\alexc\OneDrive\文件\WikiLLM\outputs\templates\bzs-report-template.html"
PDF_FILE = r"c:\Users\alexc\OneDrive\文件\WikiLLM\outputs\templates\bzs-report-template.pdf"
PPTX_FILE = r"c:\Users\alexc\OneDrive\文件\WikiLLM\outputs\templates\bzs-presentation-template.pptx"

# 2. 品牌色彩系統 (BreezySign CIS Colors)
C_PRIMARY = RGBColor(5, 120, 87)       # 翠綠 (Primary Emerald)
C_SECONDARY = RGBColor(2, 132, 199)   # 天藍 (Secondary Sky)
C_LBG = RGBColor(236, 253, 245)       # 極淺綠 (Light Background)
C_DARK = RGBColor(15, 23, 42)         # Slate 深藍灰 (Dark Text)
C_WHITE = RGBColor(255, 255, 255)
C_GRAY_LIGHT = RGBColor(241, 245, 249) # Table header background

# ---------------------------------------------------------
# A0. 產生高清 Logo 圖片 (Edge Headless 截圖技術)
# ---------------------------------------------------------
def generate_logo_pngs():
    import urllib.request
    import ssl
    from PIL import Image

    scratch_dir = r"c:\Users\alexc\OneDrive\文件\WikiLLM\scratch"
    outputs_dir = r"c:\Users\alexc\OneDrive\文件\WikiLLM\outputs"
    
    src_png = os.path.join(scratch_dir, "breezysign_logo.png")
    png_green = os.path.join(outputs_dir, "assets", "bzs-logo-green.png")
    png_white = os.path.join(outputs_dir, "assets", "bzs-logo-white.png")

    # 1. 若本地原始 PNG 不存在，優先從 brain 歷史目錄複製正確新版 Logo 圖片，防止聯網下載到官網舊資產
    if not os.path.exists(src_png):
        brain_logo = r"C:\Users\alexc\.gemini\antigravity-ide\brain\66a03437-0030-46b4-8b43-68855fce0db6\media__1779783847511.png"
        if os.path.exists(brain_logo):
            import shutil
            try:
                shutil.copy(brain_logo, src_png)
                print(f"[SUCCESS] Logo copied from brain directory: {brain_logo}")
            except Exception as e:
                print(f"[ERROR] Failed to copy logo from brain directory: {e}")
        else:
            url = "https://www.breezysign.com/breezysign_logo.png"
            print(f"Original logo not found in brain. Downloading from official website: {url}")
            ctx = ssl.create_default_context()
            ctx.check_hostname = False
            ctx.verify_mode = ssl.CERT_NONE
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
            }
            try:
                req = urllib.request.Request(url, headers=headers)
                with urllib.request.urlopen(req, context=ctx, timeout=10) as response:
                    with open(src_png, "wb") as f:
                        f.write(response.read())
                print("[SUCCESS] Successfully downloaded official logo.")
            except Exception as e_dl:
                print(f"[ERROR] Failed to download logo: {e_dl}")
                return False

    # 2. 使用 PIL 進行無損 RGBA 透明度轉換與反白處理
    try:
        img = Image.open(src_png)
        img_rgba = img.convert("RGBA")
        
        # 儲存深綠色透明版 Logo
        img_rgba.save(png_green, "PNG")
        print(f"[SUCCESS] Green logo saved to: {png_green}")
        
        # 產生純白色反白透明版 Logo (將所有非透明像素設為純白，完美保留 Alpha 抗鋸齒)
        data = img_rgba.getdata()
        new_data = []
        for item in data:
            r, g, b, a = item
            new_data.append((255, 255, 255, a))
            
        img_white = Image.new("RGBA", img_rgba.size)
        img_white.putdata(new_data)
        img_white.save(png_white, "PNG")
        print(f"[SUCCESS] White logo saved to: {png_white}")
        return True
    except Exception as e:
        print(f"[ERROR] Failed to process logo images with PIL: {e}")
        return False


# ---------------------------------------------------------
# A. PDF 轉檔功能 (自帶 Edge Headless 超時與防禦參數)
# ---------------------------------------------------------
def convert_html_to_pdf():
    edge_path = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
    if not os.path.exists(edge_path):
        print(f"Error: msedge.exe not found at {edge_path}")
        return False

    abs_html = os.path.abspath(HTML_FILE)
    abs_pdf = os.path.abspath(PDF_FILE)
    url = "file:///" + abs_html.replace("\\", "/")

    cmd = [
        edge_path,
        "--headless",
        "--disable-gpu",
        "--disable-software-rasterizer",
        "--disable-extensions",
        "--no-sandbox",
        "--print-to-pdf=" + abs_pdf,
        "--no-pdf-header-footer",
        url
    ]

    print(f"Converting {os.path.basename(abs_html)} to PDF...")
    try:
        result = subprocess.run(cmd, capture_output=True, encoding="utf-8", timeout=30)
        if result.returncode == 0:
            print(f"[SUCCESS] PDF successfully generated at: {abs_pdf}")
            print(f"File size: {os.path.getsize(abs_pdf)} bytes")
            return True
        else:
            print(f"[ERROR] Error generating PDF: {result.stderr}")
            return False
    except subprocess.TimeoutExpired:
        print("[ERROR] PDF conversion timed out (30s exceeded). Edge headless may be hung.")
        return False
    except Exception as e:
        print(f"[ERROR] An unexpected error occurred: {e}")
        return False

# ---------------------------------------------------------
# B. PPTX 簡報模板生成 (利用 python-pptx 繪製高端品牌 Slide)
# ---------------------------------------------------------
def create_pptx_template():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Helper: 設定背景顏色
    def set_slide_bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    # Helper: 繪製圓角卡片
    def add_card(slide, left, top, width, height, bg_color, border_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # Helper: 繪製實心矩形
    def add_rect(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # Helper: 新增標準內容頁標頭 (CIS 線條)
    def add_slide_header(slide, title_text):
        # 頂部翠綠線條
        add_rect(slide, 0, 0, Inches(13.33), Inches(0.15), C_PRIMARY)
        # 天藍輔助線
        add_rect(slide, 0, Inches(0.15), Inches(13.33), Inches(0.04), C_SECONDARY)
        
        # 標題文字區 (為了預防標題與右上角 Logo 重疊，限制標題寬度為 8.5 英吋)
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.5), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY

        # 右上角置入官方 Logo 圖片 (白底綠字 Logo，維持完美的 5.02:1 官方長寬比)
        logo_green_path = os.path.join(os.path.dirname(os.path.dirname(PPTX_FILE)), "assets", "bzs-logo-green.png")
        if os.path.exists(logo_green_path):
            slide.shapes.add_picture(logo_green_path, Inches(10.5), Inches(0.38), width=Inches(2.0), height=Inches(0.4))
        return title_box

    # ==========================================
    # Slide 1: 封面頁 (Cover Slide) - 品牌翠綠大氣風
    # ==========================================
    slide_cover = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_cover, C_PRIMARY)

    # 頂部天藍裝飾塊
    add_rect(slide_cover, Inches(0), Inches(0), Inches(13.33), Inches(0.2), C_SECONDARY)

    # 3. 插入官方高保真反白 Logo 圖片 (綠底白字 Logo，維持完美的 5.02:1 官方長寬比)
    logo_white_path = os.path.join(os.path.dirname(os.path.dirname(PPTX_FILE)), "assets", "bzs-logo-white.png")
    if os.path.exists(logo_white_path):
        slide_cover.shapes.add_picture(logo_white_path, Inches(1.2), Inches(1.8), width=Inches(2.61), height=Inches(0.52))

    # 封面主標題
    title_box = slide_cover.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(11.0), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "BreezySign 好好簽"
    p1.font.name = 'Microsoft JhengHei'
    p1.font.size = Pt(56)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "專屬商業演示與技術分析報告用簡報模板"
    p2.font.name = 'Microsoft JhengHei'
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_before = Pt(15)

    p3 = tf.add_paragraph()
    p3.text = "蒙恬科技 ． 電子簽章市場與技術優化小組"
    p3.font.name = 'Microsoft JhengHei'
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(200, 230, 215)
    p3.space_before = Pt(30)

    # ==========================================
    # Slide 2: 過渡頁 (Section Divider) - 科技天藍風格
    # ==========================================
    slide_div = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_div, C_SECONDARY)

    div_box = slide_div.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(10.33), Inches(3.0))
    tf_div = div_box.text_frame
    tf_div.word_wrap = True
    
    pd1 = tf_div.paragraphs[0]
    pd1.text = "PART 01"
    pd1.font.name = 'Outfit'
    pd1.font.size = Pt(32)
    pd1.font.bold = True
    pd1.font.color.rgb = RGBColor(200, 230, 255)
    
    pd2 = tf_div.add_paragraph()
    pd2.text = "地端大腦部署與開源模型選型分析"
    pd2.font.name = 'Microsoft JhengHei'
    pd2.font.size = Pt(40)
    pd2.font.bold = True
    pd2.font.color.rgb = C_WHITE
    pd2.space_before = Pt(10)

    # ==========================================
    # Slide 3: 雙欄內容頁 (Content Slide) - 經典排版
    # ==========================================
    slide_content = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_content, C_WHITE)
    add_slide_header(slide_content, "二、 雙欄式多功能資訊卡片 (經典 Light 佈局)")

    # 左側卡片 (翠綠邊框)
    add_card(slide_content, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0), C_LBG, C_PRIMARY)
    left_box = slide_content.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.6))
    tfl = left_box.text_frame
    tfl.word_wrap = True
    
    pl1 = tfl.paragraphs[0]
    pl1.text = "📌 模組 A：銷售情報與垂直行銷"
    pl1.font.name = 'Microsoft JhengHei'
    pl1.font.size = Pt(20)
    pl1.font.bold = True
    pl1.font.color.rgb = C_PRIMARY
    
    pl_body = [
        "1. 點點簽大漲 3-5 倍，改以份計費引發年費大漲輿情。",
        "2. 太平洋旅行社確定跳槽好好簽，啟用 40人吃到飽方案。",
        "3. 微型客資去重 dedupe 引擎已正式發布，防禦數據混淆。",
        "4. B2B 帳戶全面部署 tax_id 統一編號欄位，精準收束客資。"
    ]
    for item in pl_body:
        p = tfl.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(10)

    # 右側卡片 (天藍邊框)
    add_card(slide_content, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.0), RGBColor(240, 248, 255), C_SECONDARY)
    right_box = slide_content.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.0), Inches(4.6))
    tfr = right_box.text_frame
    tfr.word_wrap = True
    
    pr1 = tfr.paragraphs[0]
    pr1.text = "⚙️ 模組 B：地端算力與 RAG 技術鏈"
    pr1.font.name = 'Microsoft JhengHei'
    pr1.font.size = Pt(20)
    pr1.font.bold = True
    pr1.font.color.rgb = C_SECONDARY
    
    pr_body = [
        "1. 地端 Local LLM 優先採用 Apache 2.0 / MIT 安全開源授權。",
        "2. Qwen 2.5 7B 具備 128K context window，繁中能力頂級。",
        "3. Vector DB 採用 SQLite 驅動的 ChromaDB 實現零運維嵌入。",
        "4. LlamaIndex 框架整合 Ollama API，提供 WikiLLM 完美對接。"
    ]
    for item in pr_body:
        p = tfr.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(14)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(10)

    # ==========================================
    # Slide 4: 數據對比表格頁 (Table Slide)
    # ==========================================
    slide_table = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_table, C_WHITE)
    add_slide_header(slide_table, "三、 本土電子簽章三強 SEO / GEO 量化看板")

    # 新增表格
    rows, cols = 4, 4
    left, top, width, height = Inches(1.5), Inches(1.8), Inches(10.33), Inches(4.5)
    table_shape = slide_table.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # 設定列寬
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.53)

    # 資料填充
    table_data = [
        ["品牌名稱", "技術 SEO 評分", "GEO 能見度", "授權/安全資格"],
        ["點點簽 (DottedSign)", "80 / 100", "5.5 / 10 (中)", "AATL, ISO 27001"],
        ["律果簽 (LegalSign)", "75 / 100", "5.0 / 10 (中)", "數發部能量登錄, 憑證"],
        ["好好簽 (BreezySign) [Staging]", "85 / 100", "7.0 / 10 (高)", "能量登錄, 100% 樹狀大綱"]
    ]

    for row_idx, row in enumerate(table_data):
        for col_idx, text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            # 文字樣式與對齊
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(16)
            
            # 首行表頭 (翠綠色底)
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_PRIMARY
                p.font.bold = True
                p.font.color.rgb = C_WHITE
            # 我方數據行 (極淺綠底高亮)
            elif row_idx == 3:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LBG
                p.font.bold = True
                p.font.color.rgb = C_PRIMARY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE
                p.font.color.rgb = C_DARK

    # ==========================================
    # Slide 5: 流程圖/架構頁 (Framework Slide)
    # ==========================================
    slide_fw = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_fw, C_WHITE)
    add_slide_header(slide_fw, "四、 BreezyBrain 地端大腦非同步自動化派單流程")

    # 1. 客資輸入卡片
    add_card(slide_fw, Inches(1.0), Inches(2.2), Inches(3.2), Inches(3.5), C_LBG, C_PRIMARY)
    tb1 = slide_fw.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(3.0), Inches(3.1))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "1. 數據輸入端\n(BreezyCRM)"
    p.font.bold, p.font.size, p.font.color.rgb = True, Pt(18), C_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    p_body = ["\n名片掃描 / 外部留單", "自動 OCR 去重", "稅號 tax_id 自動匹配"]
    for b in p_body:
        pt = tf1.add_paragraph()
        pt.text = b
        pt.font.size, pt.font.color.rgb = Pt(13), C_DARK
        pt.alignment = PP_ALIGN.CENTER

    # 箭頭 1
    a1 = slide_fw.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.4), Inches(3.6), Inches(0.6), Inches(0.5))
    a1.fill.solid()
    a1.fill.fore_color.rgb = C_SECONDARY
    a1.line.fill.background()

    # 2. 地端大腦處理卡片
    add_card(slide_fw, Inches(5.1), Inches(2.2), Inches(3.2), Inches(3.5), RGBColor(240, 248, 255), C_SECONDARY)
    tb2 = slide_fw.shapes.add_textbox(Inches(5.2), Inches(2.4), Inches(3.0), Inches(3.1))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "2. 智能大腦端\n(Local LLM)"
    p.font.bold, p.font.size, p.font.color.rgb = True, Pt(18), C_SECONDARY
    p.alignment = PP_ALIGN.CENTER
    
    p_body = ["\nQwen 2.5 7B (Apache 2.0)", "RAG 範本語意匹配", "自動提取合約變數"]
    for b in p_body:
        pt = tf2.add_paragraph()
        pt.text = b
        pt.font.size, pt.font.color.rgb = Pt(13), C_DARK
        pt.alignment = PP_ALIGN.CENTER

    # 箭頭 2
    a2 = slide_fw.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.5), Inches(3.6), Inches(0.6), Inches(0.5))
    a2.fill.solid()
    a2.fill.fore_color.rgb = C_PRIMARY
    a2.line.fill.background()

    # 3. 完簽歸檔卡片
    add_card(slide_fw, Inches(9.2), Inches(2.2), Inches(3.2), Inches(3.5), C_LBG, C_PRIMARY)
    tb3 = slide_fw.shapes.add_textbox(Inches(9.3), Inches(2.4), Inches(3.0), Inches(3.1))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "3. 完簽歸檔端\n(BreezyKM)"
    p.font.bold, p.font.size, p.font.color.rgb = True, Pt(18), C_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    p_body = ["\nBreezySign API 完簽", "WikiLLM 知識庫攝入", "ChromaDB 向量化歸檔"]
    for b in p_body:
        pt = tf3.add_paragraph()
        pt.text = b
        pt.font.size, pt.font.color.rgb = Pt(13), C_DARK
        pt.alignment = PP_ALIGN.CENTER

    # 3. 保存簡報 (防禦鎖定)
    try:
        prs.save(PPTX_FILE)
        print(f"[SUCCESS] PPTX presentation generated at: {PPTX_FILE}")
        print(f"File size: {os.path.getsize(PPTX_FILE)} bytes")
    except PermissionError:
        print(f"  [WARNING] {os.path.basename(PPTX_FILE)} is locked (probably open in PowerPoint). Trying alternates...")
        saved = False
        v_idx = 2
        while not saved and v_idx < 100:
            alt_file = PPTX_FILE.replace(".pptx", f"_v{v_idx}.pptx")
            try:
                prs.save(alt_file)
                print(f"  [SUCCESS] Alternate PPTX successfully generated at: {alt_file}")
                print(f"  File size: {os.path.getsize(alt_file)} bytes")
                saved = True
            except PermissionError:
                print(f"  [WARNING] {os.path.basename(alt_file)} is also locked. Trying next...")
                v_idx += 1
        if not saved:
            print("[ERROR] Failed to save PPTX: All alternate file paths are locked.")

# ---------------------------------------------------------
# B5. 自動更新 HTML 模板中的 Base64 官方 Logo
# ---------------------------------------------------------
def update_html_logo():
    import base64
    import re
    html_path = r"c:\Users\alexc\OneDrive\文件\WikiLLM\outputs\templates\bzs-report-template.html"
    png_path = r"c:\Users\alexc\OneDrive\文件\WikiLLM\outputs\assets\bzs-logo-green.png"
    
    if not os.path.exists(html_path) or not os.path.exists(png_path):
        print(f"[ERROR] Cannot update HTML logo: files missing.")
        return False
        
    try:
        with open(png_path, "rb") as f:
            b64_str = base64.b64encode(f.read()).decode("utf-8")
            
        with open(html_path, "r", encoding="utf-8") as f:
            content = f.read()
            
        # 將舊的 SVG 向量 Logo 標籤，或者已經替換過的 img 標籤，全自動替換為最新的 Base64 PNG 標籤
        pattern_svg = r'<svg class="bzs-logo"[^>]*>.*?</svg>'
        new_logo_tag = f'<img class="bzs-logo" src="data:image/png;base64,{b64_str}" width="220" height="44" style="display:block;" alt="BreezySign">'
        
        if re.search(pattern_svg, content, flags=re.DOTALL):
            updated_content = re.sub(pattern_svg, new_logo_tag, content, flags=re.DOTALL)
        else:
            pattern_img = r'<img class="bzs-logo"[^>]*>'
            updated_content = re.sub(pattern_img, new_logo_tag, content, flags=re.DOTALL)
            
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(updated_content)
        print("[SUCCESS] HTML template logo successfully updated with Base64 PNG.")
        return True
    except Exception as e:
        print(f"[ERROR] Failed to update HTML logo: {e}")
        return False

# ---------------------------------------------------------
# C. 執行入口
# ---------------------------------------------------------
if __name__ == "__main__":
    print("Starting BZS templates generation process...")
    # 1. 產生高清 Logo 圖片
    generate_logo_pngs()

    # 1.5 自動將正版高清 Logo 以 Base64 嵌入 HTML 模板
    update_html_logo()

    # 2. 轉 PDF
    pdf_success = convert_html_to_pdf()
    
    # 3. 建立 PPTX
    create_pptx_template()
    
    print("All template generation tasks completed.")
