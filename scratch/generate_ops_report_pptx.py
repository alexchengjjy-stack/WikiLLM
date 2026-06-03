# -*- coding: utf-8 -*-
import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. 定義顏色 (對齊 BreezySign CIS)
C_PRIMARY = RGBColor(5, 120, 87)       # 翠綠 (Primary Emerald)
C_SECONDARY = RGBColor(2, 132, 199)   # 天藍 (Secondary Sky)
C_LBG = RGBColor(236, 253, 245)       # 極淺綠 (Light Background)
C_DARK = RGBColor(15, 23, 42)         # Slate 深藍灰 (Dark Text)
C_WHITE = RGBColor(255, 255, 255)
C_GRAY_LIGHT = RGBColor(241, 245, 249) # Table header background
C_RED_LIGHT = RGBColor(254, 242, 242)  # 淺紅色背景
C_RED = RGBColor(239, 68, 68)          # 紅色邊框
C_RED_DARK = RGBColor(185, 28, 28)     # 深紅色文字
C_BLUE_LIGHT = RGBColor(240, 248, 255) # 淺藍色背景

def create_report_pptx():
    outputs_dir = r"c:\Users\alexc\OneDrive\文件\WikiLLM\outputs"
    template_path = os.path.join(outputs_dir, "templates", "bzs-presentation-template.pptx")
    
    if not os.path.exists(template_path):
        print(f"[ERROR] Template PPTX not found at: {template_path}")
        return False

    # 1. 尋找最新 HTML 以對齊 timestamp 檔名
    html_pattern = os.path.join(outputs_dir, "bzs", "bzs-ops-report-*-v*.html")
    html_files = glob.glob(html_pattern)
    if not html_files:
        print("[ERROR] No HTML report files found. Cannot align timestamp.")
        return False
        
    html_files.sort(key=os.path.getmtime, reverse=True)
    latest_html = html_files[0]
    parts = os.path.basename(latest_html).split("-")
    timestamp = f"{parts[3]}-{parts[4]}" # 取得完整的 timestamp
    
    pptx_filename = f"bzs-ops-report-{timestamp}-v4.pptx"
    pptx_filepath = os.path.join(outputs_dir, "bzs", pptx_filename)

    # 2. 載入簡報 (保留格式與屬性)
    prs = Presentation(template_path)
    
    # 3. 刪除原模板簡報中所有的 mock slides (從後往前刪除)
    for i in range(len(prs.slides)-1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Helper: 設定背景顏色
    def set_slide_bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    # Helper: 繪製圓角卡片
    def add_card(slide, left, top, width, height, bg_color, border_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # Helper: 繪製實心矩形
    def add_rect(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # Helper: 新增標準內容頁標頭 (CIS 線條與 Logo)
    def add_slide_header(slide, title_text):
        # 頂部翠綠線條
        add_rect(slide, 0, 0, Inches(13.33), Inches(0.15), C_PRIMARY)
        # 天藍輔助線
        add_rect(slide, 0, Inches(0.15), Inches(13.33), Inches(0.04), C_SECONDARY)
        
        # 標題文字區
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.5), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY

        # 右上角置入官方 Logo
        logo_green_path = os.path.join(outputs_dir, "assets", "bzs-logo-green.png")
        if os.path.exists(logo_green_path):
            slide.shapes.add_picture(logo_green_path, Inches(10.5), Inches(0.38), width=Inches(2.0), height=Inches(0.4))
        return title_box

    # =========================================================
    # Slide 1: 封面頁 (Cover Slide)
    # =========================================================
    slide_cover = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_cover, C_PRIMARY)

    # 頂部天藍裝飾塊
    add_rect(slide_cover, Inches(0), Inches(0), Inches(13.33), Inches(0.2), C_SECONDARY)

    # 插入官方高保真反白 Logo
    logo_white_path = os.path.join(outputs_dir, "assets", "bzs-logo-white.png")
    if os.path.exists(logo_white_path):
        slide_cover.shapes.add_picture(logo_white_path, Inches(1.2), Inches(1.8), width=Inches(2.61), height=Inches(0.52))

    # 封面主標題
    title_box = slide_cover.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(11.0), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "BreezySign 好好簽"
    p1.font.name = 'Microsoft JhengHei'
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "2026 年 5 月營運月報與全局戰略分析"
    p2.font.name = 'Microsoft JhengHei'
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_before = Pt(15)

    p3 = tf.add_paragraph()
    p3.text = "蒙恬科技 ． 電子簽章業務與技術整合小組"
    p3.font.name = 'Microsoft JhengHei'
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(200, 230, 215)
    p3.space_before = Pt(30)

    # =========================================================
    # Slide 2: 財務營收與付費客戶結構 (雙引擎分析)
    # =========================================================
    slide_financial = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_financial, C_WHITE)
    add_slide_header(slide_financial, "一、 2026年5月財務營收與付費客戶結構")

    # 左側卡片 (營收雙引擎)
    add_card(slide_financial, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0), C_LBG, C_PRIMARY)
    left_box = slide_financial.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.6))
    tfl = left_box.text_frame
    tfl.word_wrap = True
    
    pl1 = tfl.paragraphs[0]
    pl1.text = "📊 實收總營收：NT$ 385,602"
    pl1.font.name = 'Microsoft JhengHei'
    pl1.font.size = Pt(20)
    pl1.font.bold = True
    pl1.font.color.rgb = C_PRIMARY
    
    pl_body = [
        "1. 雙引擎增長模式：",
        "  - 專案與 API 實收：NT$ 281,122 (占比 72.9%)",
        "  - SaaS 經常性訂閱：NT$ 104,480 (占比 27.1%)",
        "2. 專案營收特徵：",
        "  - 屬於一次性(One-time)系統整合款項認列",
        "  - 來自唯心醫管、百加乾燥、NX、得勝者等專案款結算",
        "  - 應與 SaaS 經常性 MRR 進行物理拆分核算"
    ]
    for item in pl_body:
        p = tfl.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 右側卡片 (SaaS 客戶結構)
    add_card(slide_financial, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.0), C_BLUE_LIGHT, C_SECONDARY)
    right_box = slide_financial.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.0), Inches(4.6))
    tfr = right_box.text_frame
    tfr.word_wrap = True
    
    pr1 = tfr.paragraphs[0]
    pr1.text = "💎 SaaS 經常性訂閱方案結構"
    pr1.font.name = 'Microsoft JhengHei'
    pr1.font.size = Pt(20)
    pr1.font.bold = True
    pr1.font.color.rgb = C_SECONDARY
    
    pr_body = [
        "1. SaaS 各方案實收明細：",
        "  - 企業方案：NT$ 73,500 (占比 70.4%)",
        "  - 專業方案：NT$ 28,800 (占比 27.6%)",
        "  - 商務方案：NT$ 1,500 (占比 1.4% - 另有大單 NT$ 13.5K 延扣)",
        "  - 加購項目(簡訊與雲端憑證)：NT$ 680",
        "2. 結構分析啟示：",
        "  - 企業方案貢獻了月實收的 83.2% (大客客單拉動)",
        "  - 專業方案家數多，為 MRR/ARR 提供穩固的續訂底座"
    ]
    for item in pr_body:
        p = tfr.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # =========================================================
    # Slide 3: 歷史營收與 MoM 趨勢 (表格)
    # =========================================================
    slide_trend = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_trend, C_WHITE)
    add_slide_header(slide_trend, "二、 2025.10 - 2026.05 SaaS 歷年實收與 MoM 趨勢")

    # 新增表格
    rows, cols = 9, 4
    left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.73), Inches(4.7)
    table_shape = slide_trend.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # 設定列寬
    table.columns[0].width = Inches(1.6)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(6.13)

    # 表格資料
    table_data = [
        ["月份", "SaaS 實收營收", "付費公司數", "MoM 增減幅度 / 結構明細"],
        ["2025-10", "NT$ 64,014", "145 家", "基準月份 (企業 $35K | 專業 $27K)"],
        ["2025-11", "NT$ 114,880", "154 家", "+79.46% (↗) (企業 $78K | 專業 $34K)"],
        ["2025-12", "NT$ 181,440", "171 家", "+57.94% (↗) (企業 $134K | 專業 $46K)"],
        ["2026-01", "NT$ 161,586", "188 家", "-10.94% (↘) (企業 $129K | 專業 $26K)"],
        ["2026-02", "NT$ 129,310", "190 家", "-19.97% (↘) (企業 $97K | 專業 $26K)"],
        ["2026-03", "NT$ 134,903", "193 家", "+4.33% (↗) (企業 $98K | 專業 $29K)"],
        ["2026-04", "NT$ 194,779", "198 家", "+44.38% (↗) (企業 $142K | 專業 $34K)"],
        ["2026-05", "NT$ 104,480", "205 家", "-46.36% (↘)* (企業 $73.5K | 商務 $1.5K* | 專業 $28.8K)"]
    ]

    for row_idx, row in enumerate(table_data):
        for col_idx, text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Microsoft JhengHei'
            
            # 第一行表頭
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_SECONDARY
                p.font.bold = True
                p.font.size = Pt(13)
                p.font.color.rgb = C_WHITE
                p.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE if row_idx % 2 != 0 else C_GRAY_LIGHT
                p.font.size = Pt(10.5)
                p.font.color.rgb = C_DARK
                if col_idx < 3:
                    p.alignment = PP_ALIGN.CENTER
                    if col_idx == 0:
                        p.font.bold = True
                else:
                    p.alignment = PP_ALIGN.LEFT
                    if "+" in text:
                        p.font.color.rgb = C_PRIMARY
                        p.font.bold = True
                    elif "-" in text and "MoM" not in text:
                        p.font.color.rgb = RGBColor(185, 28, 28)
                        p.font.bold = True

    # 底部註記文字
    note_box = slide_trend.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.73), Inches(0.5))
    tf_n = note_box.text_frame
    tf_n.word_wrap = True
    p_n = tf_n.paragraphs[0]
    p_n.text = "* 註：2026-05 SaaS 實收金流因大單（太平洋旅行社 $60K）與商務方案大單（$13.5K）生效扣款期延遲為 6/1 呈技術下降。若併計當月專案實收 NT$ 281,122，5 月實收總營收達 NT$ 385,602，總體實收 MoM 實際為 +97.97% 強勁成長。"
    p_n.font.name = 'Microsoft JhengHei'
    p_n.font.size = Pt(9.5)
    p_n.font.italic = True
    p_n.font.color.rgb = RGBColor(100, 110, 120)

    # =========================================================
    # Slide 4: 新增註冊與獲客漏斗成效 (SaaS 漏斗與 CPA 分析)
    # =========================================================
    slide_funnel = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_funnel, C_WHITE)
    add_slide_header(slide_funnel, "三、 新增註冊、獲客漏斗與 CPA 獲客成效")

    # 左側卡片 (獲客漏斗與 CPL 指標)
    add_card(slide_funnel, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0), C_LBG, C_PRIMARY)
    left_box_f = slide_funnel.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.6))
    tfl_f = left_box_f.text_frame
    tfl_f.word_wrap = True
    
    plf1 = tfl_f.paragraphs[0]
    plf1.text = "🚀 獲客漏斗與 CPA 獲客成效"
    plf1.font.name = 'Microsoft JhengHei'
    plf1.font.size = Pt(20)
    plf1.font.bold = True
    plf1.font.color.rgb = C_PRIMARY
    
    plf_body = [
        "1. 月度獲客指標明細：",
        "  - 新增註冊公司數：312 家 (月免費註冊 base 擴大)",
        "  - 當月累計電訪註冊 Leads：30 家",
        "  - 表達「有興趣」：15 家 (轉換率高達 50%)",
        "  - 歸入「較高意願」：9 家，由 CSM 持續專人輔導",
        "2. 廣告費用與 CPL 精算 (當月預算 NT$ 145,080)：",
        "  - 寬口徑 CPL (按註冊計)：NT$ 465",
        "  - 窄口徑 CPL (按高意願計)：NT$ 1,792"
    ]
    for item in plf_body:
        p = tfl_f.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 右側卡片 (業界基準比較)
    add_card(slide_funnel, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.0), C_BLUE_LIGHT, C_SECONDARY)
    right_box_f = slide_funnel.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.0), Inches(4.6))
    tfr_f = right_box_f.text_frame
    tfr_f.word_wrap = True
    
    prf1 = tfr_f.paragraphs[0]
    prf1.text = "💡 SaaS 核心指標釋義與業界基準"
    prf1.font.name = 'Microsoft JhengHei'
    prf1.font.size = Pt(20)
    prf1.font.bold = True
    prf1.font.color.rgb = C_SECONDARY
    
    prf_body = [
        "1. B2B 電子簽章獲客成本對照：",
        "  - 台灣 B2B CPA 均值：NT$ 1,000 ~ 3,200",
        "  - 全球 B2B CPA 均值：NT$ 2,500 ~ 6,400+",
        "  - 好好簽寬口徑 (NT$ 465) 極具競爭優勢",
        "  - 窄口徑 (NT$ 1,792) 落在台灣正常均值內，成效極佳",
        "2. LTV : CAC 財務健康度：",
        "  - 企業年約客單價 LTV 估計達 NT$ 120,000",
        "  - 窄口徑 LTV:CAC 獲客效率高達 67 倍，利潤空間大"
    ]
    for item in prf_body:
        p = tfr_f.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(6)

    # =========================================================
    # Slide 5: SaaS 歷年四大維度與指標演進 (2024-2026) -> [NEW]
    # =========================================================
    slide_dimensions = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_dimensions, C_WHITE)
    add_slide_header(slide_dimensions, "四、 SaaS 歷年四大維度與指標演進 (2024-2026)")

    # 4 個小卡片 (2x2)
    # 卡片 1 (獲客)
    add_card(slide_dimensions, Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.3), C_LBG, C_PRIMARY)
    box_d1 = slide_dimensions.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(2.1))
    tf_d1 = box_d1.text_frame
    tf_d1.word_wrap = True
    pd1 = tf_d1.paragraphs[0]
    pd1.text = "🎯 獲客維度 (Acquisition)"
    pd1.font.name = 'Microsoft JhengHei'
    pd1.font.size = Pt(16)
    pd1.font.bold = True
    pd1.font.color.rgb = C_PRIMARY
    
    d1_items = [
        "• 歷年註冊：2024下半年 2,126 次 -> 2025年 3,738 次 -> 2026年前五月 1,620 次，累計用戶基數快速擴大。",
        "• 5月 Ads 花費 NT$ 145,080，B2B 窄口徑 CPA 為 NT$ 1,792。"
    ]
    for item in d1_items:
        p = tf_d1.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(10.5)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(3)

    # 卡片 2 (營收)
    add_card(slide_dimensions, Inches(6.9), Inches(1.5), Inches(5.6), Inches(2.3), C_BLUE_LIGHT, C_SECONDARY)
    box_d2 = slide_dimensions.shapes.add_textbox(Inches(7.1), Inches(1.6), Inches(5.2), Inches(2.1))
    tf_d2 = box_d2.text_frame
    tf_d2.word_wrap = True
    pd2 = tf_d2.paragraphs[0]
    pd2.text = "💰 營收維度 (Revenue)"
    pd2.font.name = 'Microsoft JhengHei'
    pd2.font.size = Pt(16)
    pd2.font.bold = True
    pd2.font.color.rgb = C_SECONDARY
    
    d2_items = [
        "• 經常性實收：2024下半年 NT$ 95.4K -> 2025年 NT$ 1.26M -> 2026前五個月累計達 NT$ 728.7K。",
        "• 5月單月實收創 NT$ 385,602 歷史新高 (SaaS 104K + 專案 281K)。"
    ]
    for item in d2_items:
        p = tf_d2.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(10.5)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(3)

    # 卡片 3 (留存)
    add_card(slide_dimensions, Inches(0.8), Inches(4.1), Inches(5.6), Inches(2.3), C_BLUE_LIGHT, C_SECONDARY)
    box_d3 = slide_dimensions.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(5.2), Inches(2.1))
    tf_d3 = box_d3.text_frame
    tf_d3.word_wrap = True
    pd3 = tf_d3.paragraphs[0]
    pd3.text = "🔄 留存維度 (Retention)"
    pd3.font.name = 'Microsoft JhengHei'
    pd3.font.size = Pt(16)
    pd3.font.bold = True
    pd3.font.color.rgb = C_SECONDARY
    
    d3_items = [
        "• 留存防線：利用未公開之「商務方案 (降級價)」作為 Down-selling 籌碼，防守客戶流失。",
        "• 透過 PipeDrive 記錄「任務活躍度」與 Line 諮詢，建立免費版活躍度預警，控制年流失率在 5% 以內。"
    ]
    for item in d3_items:
        p = tf_d3.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(10.5)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(3)

    # 卡片 4 (價值)
    add_card(slide_dimensions, Inches(6.9), Inches(4.1), Inches(5.6), Inches(2.3), C_LBG, C_PRIMARY)
    box_d4 = slide_dimensions.shapes.add_textbox(Inches(7.1), Inches(4.2), Inches(5.2), Inches(2.1))
    tf_d4 = box_d4.text_frame
    tf_d4.word_wrap = True
    pd4 = tf_d4.paragraphs[0]
    pd4.text = "📈 價值維度 (Value)"
    pd4.font.name = 'Microsoft JhengHei'
    pd4.font.size = Pt(16)
    pd4.font.bold = True
    pd4.font.color.rgb = C_PRIMARY
    
    d4_items = [
        "• LTV 估計值：以流失率 5% 估算，企業年約客單價帶來的終身價值 (LTV) 高達 NT$ 120,000。",
        "• 獲客效率：LTV:CAC 達 67 倍，投資回報率極高。2026上半年實收已完全收回獲客成本 (回本期 < 12 個月)。"
    ]
    for item in d4_items:
        p = tf_d4.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(10.5)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(3)

    # =========================================================
    # Slide 6: 企業客戶畫像分析 (一)：行業深耕客群 -> [NEW]
    # =========================================================
    slide_personas1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_personas1, C_WHITE)
    add_slide_header(slide_personas1, "五、 企業客戶畫像分析 (一)：垂直深耕客群")

    # 3 個垂直卡片 (3 Columns)
    # 卡片 1 (醫療)
    add_card(slide_personas1, Inches(0.8), Inches(1.5), Inches(3.6), Inches(5.0), C_LBG, C_PRIMARY)
    box_p1 = slide_personas1.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.2), Inches(4.6))
    tf_p1 = box_p1.text_frame
    tf_p1.word_wrap = True
    pp1 = tf_p1.paragraphs[0]
    pp1.text = "🏥 醫療院所與生技保健"
    pp1.font.name = 'Microsoft JhengHei'
    pp1.font.size = Pt(18)
    pp1.font.bold = True
    pp1.font.color.rgb = C_PRIMARY
    
    p1_items = [
        "• 代表：得勝者(盧森眼科)、衡星牙醫、膚適美。",
        "• 特徵與痛點：強烈要求 API 座標完簽、PDF 轉 Dicom 回寫，及離線暫存與 NTP 3天內校時的法規需求。",
        "• 渠道特點：通訊偏向使用 LINE 或簡訊簽署，而非 Email。"
    ]
    for item in p1_items:
        p = tf_p1.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(12)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 卡片 2 (轉單)
    add_card(slide_personas1, Inches(4.8), Inches(1.5), Inches(3.6), Inches(5.0), C_BLUE_LIGHT, C_SECONDARY)
    box_p2 = slide_personas1.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.2), Inches(4.6))
    tf_p2 = box_p2.text_frame
    tf_p2.word_wrap = True
    pp2 = tf_p2.paragraphs[0]
    pp2.text = "🔥 點點簽漲價轉單大戶"
    pp2.font.name = 'Microsoft JhengHei'
    pp2.font.size = Pt(18)
    pp2.font.bold = True
    pp2.font.color.rgb = C_SECONDARY
    
    p2_items = [
        "• 代表：太平洋旅行社、福安管理顧問、海沃管顧。",
        "• 特徵與痛點：因點點簽改為「按件計費」導致預算大漲 3-5 倍，轉投我方「吃到飽年租方案」。",
        "• 核心訴求：要求 UNIFY 範本權限集中控制 (限制子帳號自建)，及數發部能量登錄許可證書。"
    ]
    for item in p2_items:
        p = tf_p2.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(12)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 卡片 3 (不動產)
    add_card(slide_personas1, Inches(8.8), Inches(1.5), Inches(3.6), Inches(5.0), C_LBG, C_PRIMARY)
    box_p3 = slide_personas1.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.2), Inches(4.6))
    tf_p3 = box_p3.text_frame
    tf_p3.word_wrap = True
    pp3 = tf_p3.paragraphs[0]
    pp3.text = "🏠 不動產租售與代銷"
    pp3.font.name = 'Microsoft JhengHei'
    pp3.font.size = Pt(18)
    pp3.font.bold = True
    pp3.font.color.rgb = C_PRIMARY
    
    p3_items = [
        "• 代表：第一建經、佶星廣告、二房東星鴻。",
        "• 特徵與痛點：業務現場看房時需平板面簽，高法律防禦需求，避免客戶賴帳。",
        "• 核心訴求：要求「邊簽名邊錄影」聲明錄影防賴與 AATL 數位憑證，強化法庭上之證據力。"
    ]
    for item in p3_items:
        p = tf_p3.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(12)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # =========================================================
    # Slide 7: 企業客戶畫像分析 (二)：批量與地端客群 -> [NEW]
    # =========================================================
    slide_personas2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_personas2, C_WHITE)
    add_slide_header(slide_personas2, "六、 企業客戶畫像分析 (二)：批量與地端客群")

    # 3 個垂直卡片 (3 Columns)
    # 卡片 1 (補教)
    add_card(slide_personas2, Inches(0.8), Inches(1.5), Inches(3.6), Inches(5.0), C_BLUE_LIGHT, C_SECONDARY)
    box_p4 = slide_personas2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.2), Inches(4.6))
    tf_p4 = box_p4.text_frame
    tf_p4.word_wrap = True
    pp4 = tf_p4.paragraphs[0]
    pp4.text = "📚 大型教育與補教機構"
    pp4.font.name = 'Microsoft JhengHei'
    pp4.font.size = Pt(18)
    pp4.font.bold = True
    pp4.font.color.rgb = C_SECONDARY
    
    p4_items = [
        "• 代表：智基科技 (志光公職)。",
        "• 特徵與痛點：年簽約量 3-5 萬份，決策鏈長，要求極具競爭力的階梯式報價 (Tiered Pricing)。",
        "• 核心訴求：強烈依賴「公開表單簽」與批量發送，降低大批量發送的行政人力成本與錯誤率。"
    ]
    for item in p4_items:
        p = tf_p4.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(12)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 卡片 2 (地端)
    add_card(slide_personas2, Inches(4.8), Inches(1.5), Inches(3.6), Inches(5.0), C_LBG, C_PRIMARY)
    box_p5 = slide_personas2.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.2), Inches(4.6))
    tf_p5 = box_p5.text_frame
    tf_p5.word_wrap = True
    pp5 = tf_p5.paragraphs[0]
    pp5.text = "🏢 企業集團與數位行銷"
    pp5.font.name = 'Microsoft JhengHei'
    pp5.font.size = Pt(18)
    pp5.font.bold = True
    pp5.font.color.rgb = C_PRIMARY
    
    p5_items = [
        "• 代表：台北 101 地端 BPM 部署專案。",
        "• 特徵與痛點：要求嚴格的企業安全內控，Docker 地端部署，提供原始代碼安裝包，放寬內部 WAF 規則。",
        "• 核心訴求：系統異常需串接內部監控 (isHealth API)，支持雲端憑證發送與多種組織審批流。"
    ]
    for item in p5_items:
        p = tf_p5.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(12)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 卡片 3 (防禦)
    add_card(slide_personas2, Inches(8.8), Inches(1.5), Inches(3.6), Inches(5.0), C_RED_LIGHT, C_RED)
    box_p6 = slide_personas2.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.2), Inches(4.6))
    tf_p6 = box_p6.text_frame
    tf_p6.word_wrap = True
    pp6 = tf_p6.paragraphs[0]
    pp6.text = "⚠️ 技術利潤防禦邊界"
    pp6.font.name = 'Microsoft JhengHei'
    pp6.font.size = Pt(18)
    pp6.font.bold = True
    pp6.font.color.rgb = C_RED_DARK
    
    p6_items = [
        "• 代表：聖美麗健康管理 (主動婉拒)。",
        "• 特徵與痛點：PDF 檔案過大 (健檢影像檔 > 28MB)，嵌入 AATL 憑證時易造成伺服器負載超時簽署失敗。",
        "• 決策：為此客製會破壞毛利，主動婉拒年約，確立「單檔 10MB 憑證限制」的技術防衛邊界。"
    ]
    for item in p6_items:
        p = tf_p6.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(12)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # =========================================================
    # Slide 8: 客戶成功協助新購轉化案例 (CSM Onboarding 案例)
    # =========================================================
    slide_onboarding = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_onboarding, C_WHITE)
    add_slide_header(slide_onboarding, "七、 客戶成功主動協助新購轉化案例")

    # 左側卡片 (豐盛富足與自強基金會)
    add_card(slide_onboarding, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0), C_LBG, C_PRIMARY)
    left_box_o = slide_onboarding.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.6))
    tfl_o = left_box_o.text_frame
    tfl_o.word_wrap = True
    
    plo1 = tfl_o.paragraphs[0]
    plo1.text = "🏢 企業方案 Onboarding 成功案例"
    plo1.font.name = 'Microsoft JhengHei'
    plo1.font.size = Pt(20)
    plo1.font.bold = True
    plo1.font.color.rgb = C_PRIMARY
    
    plo_body = [
        "1. 豐盛富足資產管理 (企業方案月約)：",
        "  - 客戶註冊後因流程冷啟動，未使用簽署",
        "  - CSM 主動電訪關懷，提供一對一教學",
        "  - 引導建立常用合約範本，於 5/7 順利訂閱企業方案",
        "2. 自強工業科學基金會 (企業方案短期約)：",
        "  - 註冊後有多人協作與權限組織管理需求",
        "  - CSM 團隊迅速排定線上系統 Demo 展示",
        "  - 協助配置組織架構，於 5/15 成交 2 個月企業租約"
    ]
    for item in plo_body:
        p = tfl_o.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 右側卡片 (富友旅行社成功案例)
    add_card(slide_onboarding, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.0), C_BLUE_LIGHT, C_SECONDARY)
    right_box_o = slide_onboarding.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.0), Inches(4.6))
    tfr_o = right_box_o.text_frame
    tfr_o.word_wrap = True
    
    pro1 = tfr_o.paragraphs[0]
    pro1.text = "✈️ 專業方案年約新簽成功案例"
    pro1.font.name = 'Microsoft JhengHei'
    pro1.font.size = Pt(20)
    pro1.font.bold = True
    pro1.font.color.rgb = C_SECONDARY
    
    pro_body = [
        "1. 富友旅行社 (專業方案年約)：",
        "  - 旅行社面臨「代收轉付收據」與合約大量傳簽痛點",
        "  - CSM 針對旅遊業痛點客製化傳簽流程引導",
        "  - 協助解決線上付款與多個承辦帳號授權配置",
        "  - 客戶滿意整體服務，於 5/8 訂閱專業方案年約",
        "2. 啟示：",
        "  - 超過 50% 註冊 Leads 存在冷啟動抗性",
        "  - CSM 主動電訪介入是提升 PLG 轉化率的關鍵錨點"
    ]
    for item in pro_body:
        p = tfr_o.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(6)

    # =========================================================
    # Slide 9: 退訂與流失結案歷程 (主動退守防線)
    # =========================================================
    slide_churn = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_churn, C_WHITE)
    add_slide_header(slide_churn, "八、 退訂流失與專案結案歷程分析")

    # 左側卡片 (聖美麗健健大檔案限制)
    add_card(slide_churn, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0), C_RED_LIGHT, C_RED)
    left_box_c = slide_churn.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.6))
    tfl_c = left_box_c.text_frame
    tfl_c.word_wrap = True
    
    plc1 = tfl_c.paragraphs[0]
    plc1.text = "⚠️ 聖美麗健檢大檔案限制主動婉拒"
    plc1.font.name = 'Microsoft JhengHei'
    plc1.font.size = Pt(20)
    plc1.font.bold = True
    plc1.font.color.rgb = C_RED_DARK
    
    plc_body = [
        "1. 客戶背景與需求特徵：",
        "  - 聖美麗健檢報告富含影像，單一 PDF 檔高達 28MB",
        "2. 技術瓶頸與防衛限制：",
        "  - 超大檔案在嵌入 AATL 數位憑證時易造成伺服器負載超時，進而簽署失敗",
        "  - 為其修改底層系統架構將嚴重破壞好好簽的毛利",
        "3. 決策與結案歷程：",
        "  - CSM 與技術小組評估後，主動予以婉拒年約",
        "  - 本案正式確立我方「單檔 10MB 與 AATL」防禦邊界"
    ]
    for item in plc_body:
        p = tfl_c.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 右側卡片 (恩主公醫院臨櫃專案)
    add_card(slide_churn, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.0), C_GRAY_LIGHT, C_DARK)
    right_box_c = slide_churn.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.0), Inches(4.6))
    tfr_c = right_box_c.text_frame
    tfr_c.word_wrap = True
    
    prc1 = tfr_c.paragraphs[0]
    prc1.text = "📋 恩主公醫院臨櫃專案結案分析"
    prc1.font.name = 'Microsoft JhengHei'
    prc1.font.size = Pt(20)
    prc1.font.bold = True
    prc1.font.color.rgb = C_DARK
    
    prc_body = [
        "1. 專案背景與評估範圍：",
        "  - 評估設置電子簽章平台供院內民眾現場臨櫃簽署",
        "2. 結案原因說明：",
        "  - 5/20 院方正式來信，因今年院內資源與預算配置已滿，無法支援此案開發聯絡，本案正式結案",
        "3. 經驗與策略啟示：",
        "  - 公立/大型醫療院所 IT 決策週期長，且易受預算排擠",
        "  - 業務與技術資源優先專注於中小型自費/連鎖診所"
    ]
    for item in prc_body:
        p = tfr_c.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # =========================================================
    # Slide 10: 重大專案與 API 技術串接進程 (一)
    # =========================================================
    slide_proj1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_proj1, C_WHITE)
    add_slide_header(slide_proj1, "九、 重大專案與 API 技術串接進程 (一)")

    # 新增表格
    rows, cols = 4, 3
    left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.0)
    table_shape1 = slide_proj1.shapes.add_table(rows, cols, left, top, width, height)
    table1 = table_shape1.table

    table1.columns[0].width = Inches(2.5)
    table1.columns[1].width = Inches(3.2)
    table1.columns[2].width = Inches(6.03)

    table_data1 = [
        ["合作夥伴 / 專案", "本月進展 / 里程碑", "技術對接與合規說明"],
        ["鼎新電腦 (ISV 合作)", "API 串接完成與調優", "完成串接並優化 GCP 載入慢問題。已將 API 拋出連結有效時效由 60 秒調增為 15 分鐘。全力備戰 6/11 直播。"],
        ["百加資通 (BPM 合作)", "專案款實收開立發票", "偉勝乾燥專案實收 NT$ 60,000 已於 5/30 開立發票；巨虹實收 NT$ 38,225。凌越生醫已正式結案。"],
        ["中華電信 (合規保障)", "瑞飛智慧憑證合規澄清", "釐清好好簽 AATL 憑證確實由中華電信 CA 提供。對齊中華電信 LTV 完整時戳，確保推定法律效力。"]
    ]

    for row_idx, row in enumerate(table_data1):
        for col_idx, text in enumerate(row):
            cell = table1.cell(row_idx, col_idx)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Microsoft JhengHei'
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_PRIMARY
                p.font.bold = True
                p.font.size = Pt(14)
                p.font.color.rgb = C_WHITE
                p.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE if row_idx % 2 != 0 else C_GRAY_LIGHT
                p.font.size = Pt(12)
                p.font.color.rgb = C_DARK
                if col_idx < 2:
                    p.alignment = PP_ALIGN.CENTER
                    p.font.bold = True
                else:
                    p.alignment = PP_ALIGN.LEFT

    # =========================================================
    # Slide 11: 重大專案與 API 技術串接進程 (二)
    # =========================================================
    slide_proj2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_proj2, C_WHITE)
    add_slide_header(slide_proj2, "十、 重大專案與 API 技術串接進程 (二)")

    # 新增表格
    rows, cols = 4, 3
    left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.0)
    table_shape2 = slide_proj2.shapes.add_table(rows, cols, left, top, width, height)
    table2 = table_shape2.table

    table2.columns[0].width = Inches(2.5)
    table2.columns[1].width = Inches(3.2)
    table2.columns[2].width = Inches(6.03)

    table_data2 = [
        ["合作夥伴 / 專案", "本月進展 / 里程碑", "技術對接與合規說明"],
        ["得勝者 (HIS/PACS)", "旗下診所 7 月上線準備", "盧森與東港盧森兩案共實收 NT$ 50,000。已設計規避斷線合規的「地端離線暫存與中華電信 NTP 校時機制」。洽談商之器 PACS 後台電簽。"],
        ["台北 101 (地端專案)", "獨立地端專案交付", "此為獨立專案 (非中華電信引薦)。已交付 HiCloud 與 DMZ 安裝包及地端原始碼。正討論 isHealth API 檢測機制。"],
        ["太平洋旅行社 (大單)", "40人企業方案開通", "電匯 NT$ 60,000 入帳 (5月實收)，UNIFY 權限共享範本配置完成，正式開通 (6/1 生效)。"]
    ]

    for row_idx, row in enumerate(table_data2):
        for col_idx, text in enumerate(row):
            cell = table2.cell(row_idx, col_idx)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Microsoft JhengHei'
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_SECONDARY
                p.font.bold = True
                p.font.size = Pt(14)
                p.font.color.rgb = C_WHITE
                p.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE if row_idx % 2 != 0 else C_GRAY_LIGHT
                p.font.size = Pt(12)
                p.font.color.rgb = C_DARK
                if col_idx < 2:
                    p.alignment = PP_ALIGN.CENTER
                    p.font.bold = True
                else:
                    p.alignment = PP_ALIGN.LEFT

    # =========================================================
    # Slide 12: 潛力進行中專案與 B2B2C 合作進程
    # =========================================================
    slide_pipeline = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_pipeline, C_WHITE)
    add_slide_header(slide_pipeline, "十一、 潛力進行中專案與 B2B2C 合作進程")

    # 左側卡片 (通路合作與 API 大單)
    add_card(slide_pipeline, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0), C_LBG, C_PRIMARY)
    left_box_p = slide_pipeline.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.6))
    tfl_p = left_box_p.text_frame
    tfl_p.word_wrap = True
    
    plp1 = tfl_p.paragraphs[0]
    plp1.text = "🤝 共享辦公合作與 API 轉單大單"
    plp1.font.name = 'Microsoft JhengHei'
    plp1.font.size = Pt(20)
    plp1.font.bold = True
    plp1.font.color.rgb = C_PRIMARY
    
    plp_body = [
        "1. 大瀚 GTB 租客訂閱合作案 (B2B2C)：",
        "  - 專屬合作體驗 Landing Page 設計確認",
        "  - 預計 6 月初上線，提供大瀚租客 3 個月免費體驗",
        "  - 共享辦公生態系獲取租客訂閱之增長模型",
        "2. 福安健康與職安 API (轉單大戶)：",
        "  - 年簽署量 2 萬份，API 專案報價 12 萬簽約中",
        "  - 我方已提供 ISO27001 與數發部能量登錄證書",
        "3. 聖洋科技 (cacafly) API 串接：",
        "  - 年簽署量 8K ~ 10K 份，正進行多品牌動態 Logo 規格規劃，初步報價每份 NT$ 20"
    ]
    for item in plp_body:
        p = tfl_p.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(12.5)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(6)

    # 右側卡片 (其他跟進中案源)
    add_card(slide_pipeline, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.0), C_BLUE_LIGHT, C_SECONDARY)
    right_box_p = slide_pipeline.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.0), Inches(4.6))
    tfr_p = right_box_p.text_frame
    tfr_p.word_wrap = True
    
    prp1 = tfr_p.paragraphs[0]
    prp1.text = "📈 進行中潛力轉單案源跟進"
    prp1.font.name = 'Microsoft JhengHei'
    prp1.font.size = Pt(20)
    prp1.font.bold = True
    prp1.font.color.rgb = C_SECONDARY
    
    prp_body = [
        "1. 神坊資訊 (小樹購 API)：預估年用量 2,000 份，正提供金融業實績，議價年費中",
        "2. 聯合線上 (udn) API：串接完成，進行流程防偽測試中",
        "3. 鴻運聯邦 (汽車回收)：月用量 400-500 份，測試體驗中",
        "4. 海沃管理顧問 (點點簽受災戶)：年用量 200-300 份，體驗版試用至 6/10",
        "5. 麻吉行得通：500-600 份年約，對接 8/3 到期轉換",
        "6. 耐斯旅行社：定型化契約 Line 傳簽測試通過，引導月費制訂閱中",
        "7. 星鴻高雄二房東：爭取二房東自營品牌全面導入"
    ]
    for item in prp_body:
        p = tfr_p.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(5)

    # =========================================================
    # Slide 13: 電子簽章能量登錄競品情報普查快照 (2026年6月) -> [NEW]
    # =========================================================
    slide_competitors = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_competitors, C_WHITE)
    add_slide_header(slide_competitors, "十二、 電子簽章能量登錄與競品情報普查快照")

    # 新增普查表格
    rows, cols = 5, 4
    left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.73), Inches(4.7)
    table_shape3 = slide_competitors.shapes.add_table(rows, cols, left, top, width, height)
    table3 = table_shape3.table

    # 設定列寬
    table3.columns[0].width = Inches(1.6)
    table3.columns[1].width = Inches(2.3)
    table3.columns[2].width = Inches(5.33)
    table3.columns[3].width = Inches(2.5)

    # 表格資料
    table_data3 = [
        ["競品名稱", "能量登錄證號與有效期", "產品與行銷最新動態 (100% 正式站普查)", "核心弱點與痛點"],
        ["好好簽 (我方)", "113電簽0008\n(至 116/08/13)", "登錄證書已正式上線官網。主打不限件數吃到飽年約及 LINE 傳簽/聲明錄影防賴。", "單檔 10MB 與 AATL 技術限制 (技術防禦邊界，確保毛利)"],
        ["點點簽\n(DottedSign)", "113電簽0003\n(至 115/08/29)", "自 2026-04-21 強制廢止舊企業方案續訂，改按件計費 ($45-50/件)。整合 BizForm、MCP 大模型語意傳簽。", "價格暴增 3-5 倍，用戶抗性極大，促成大量中大型客戶流失與跳槽。"],
        ["律果簽\n(LegalSign)", "113-電簽-0005\n(至 116/01/15)", "主推 AI 法務助理「法樂多」，主打 30 秒自動審約、CLM 合約生命週期管理。", "大批量傳簽時系統 Loading 過慢 (旅行社測試需 10 分鐘) 為其硬傷。"],
        ["FastSIGN\n(全景)", "113-電簽-0001\n(至 115/07/21)", "推出 IDExpert Cloud 零信任認證，主打 PQC 後量子密碼學遷移與三階段驗證。", "線上多媒體行銷較少，GEO 能見度與市場聲音偏低。"]
    ]

    for row_idx, row in enumerate(table_data3):
        for col_idx, text in enumerate(row):
            cell = table3.cell(row_idx, col_idx)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Microsoft JhengHei'
            
            # 第一行表頭
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_PRIMARY
                p.font.bold = True
                p.font.size = Pt(12)
                p.font.color.rgb = C_WHITE
                p.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE if row_idx % 2 != 0 else C_GRAY_LIGHT
                p.font.size = Pt(10.5)
                p.font.color.rgb = C_DARK
                if col_idx < 2:
                    p.alignment = PP_ALIGN.CENTER
                    if col_idx == 0:
                        p.font.bold = True
                else:
                    p.alignment = PP_ALIGN.LEFT

    # =========================================================
    # Slide 14: 業務前線反駁話術 Battle Cards -> [NEW]
    # =========================================================
    slide_battlecards = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_battlecards, C_WHITE)
    add_slide_header(slide_battlecards, "十三、 業務前線反駁話術 Battle Cards")

    # 左側卡片 (點點簽對抗卡)
    add_card(slide_battlecards, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0), C_LBG, C_PRIMARY)
    left_box_b = slide_battlecards.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.6))
    tfl_b = left_box_b.text_frame
    tfl_b.word_wrap = True
    
    plb1 = tfl_b.paragraphs[0]
    plb1.text = "⚔️ 點點簽 (DottedSign) 對抗話術"
    plb1.font.name = 'Microsoft JhengHei'
    plb1.font.size = Pt(20)
    plb1.font.bold = True
    plb1.font.color.rgb = C_PRIMARY
    
    plb_body = [
        "1. 客戶提及：點點簽功能多，且整合了許多大型系統。",
        "2. 前線反駁與話術錨點：",
        "  - 報價分析：「點點簽自 2026/04/21 起強制廢止了舊版企業方案的續約，全面改按件計費。到期後成本將暴增 3-5 倍，預算壓力極大。」",
        "  - 性價比：「好好簽提供『吃到飽不限件數年租』方案，在中大用量客戶群中具備壓倒性的預算優勢。」",
        "  - 合規保障：「我方具備國家數發部核可《113電簽0008》證書，且提供獨家 LINE 傳簽與現場錄影防賴，背靠蒙恬科技，效力與保障更佳。」"
    ]
    for item in plb_body:
        p = tfl_b.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 右側卡片 (律果簽對抗卡)
    add_card(slide_battlecards, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.0), C_BLUE_LIGHT, C_SECONDARY)
    right_box_b = slide_battlecards.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.0), Inches(4.6))
    tfr_b = right_box_b.text_frame
    tfr_b.word_wrap = True
    
    prb1 = tfr_b.paragraphs[0]
    prb1.text = "⚔️ 律果簽 (LegalSign) 對抗話術"
    prb1.font.name = 'Microsoft JhengHei'
    prb1.font.size = Pt(20)
    prb1.font.bold = True
    prb1.font.color.rgb = C_SECONDARY
    
    prb_body = [
        "1. 客戶提及：律果簽提供 AI 審約法樂多，法務功能更豐富。",
        "2. 前線反駁與話術錨點：",
        "  - 效率硬傷：「律果簽主打 AI 審約與 CLM，但在高頻大批量傳簽時系統加載極慢，旅行社大戶測試時需加載 10 分鐘，嚴重阻礙前線成交時效。」",
        "  - 專注本質：「好好簽採用高性能簽署引擎，流暢穩定，提供實在、極速、法庭證據力強的完簽體驗。」",
        "  - 合規認證：「好好簽已通過國家數發部能量登錄核可，提供最高規格的中華電信 AATL 數位憑證，法律效力無懈可擊。」"
    ]
    for item in prb_body:
        p = tfr_b.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # =========================================================
    # Slide 15: 點點簽轉單效應與後續執行建議
    # =========================================================
    slide_churn_eff = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_churn_eff, C_WHITE)
    add_slide_header(slide_churn_eff, "十四、 點點簽轉單效應與後續執行建議")

    # 左側卡片 (點點簽轉單效應)
    add_card(slide_churn_eff, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0), C_LBG, C_PRIMARY)
    left_box_e = slide_churn_eff.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.6))
    tfl_e = left_box_e.text_frame
    tfl_e.word_wrap = True
    
    ple1 = tfl_e.paragraphs[0]
    ple1.text = "🔥 點點簽漲價與轉單效應分析"
    ple1.font.name = 'Microsoft JhengHei'
    ple1.font.size = Pt(20)
    ple1.font.bold = True
    ple1.font.color.rgb = C_PRIMARY
    
    ple_body = [
        "1. 計費模型更動：",
        "  - 改為以件計費 (單份約 NT$45-50)",
        "  - 續約報價大增 3-5 倍，用戶預算抗性大",
        "2. 好好簽定定價優勢：",
        "  - 堅持「吃到飽年租」模式，在中大用量客戶群中具備壓倒性性價比",
        "3. 轉單實績與商機：",
        "  - 太平洋旅行社 (2k份/年) 已付款轉單",
        "  - 福安健康 (2萬份/年) API 積極簽約中",
        "  - 海沃、麻吉行得通等轉單積極跟進中"
    ]
    for item in ple_body:
        p = tfl_e.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(13)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(6)

    # 右側卡片 (營運建議)
    add_card(slide_churn_eff, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.0), C_BLUE_LIGHT, C_SECONDARY)
    right_box_e = slide_churn_eff.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.0), Inches(4.6))
    tfr_e = right_box_e.text_frame
    tfr_e.word_wrap = True
    
    pre1 = tfr_e.paragraphs[0]
    pre1.text = "📋 團隊後續重點執行建議"
    pre1.font.name = 'Microsoft JhengHei'
    pre1.font.size = Pt(20)
    pre1.font.bold = True
    pre1.font.color.rgb = C_SECONDARY
    
    pre_body = [
        "1. 建立財務與 CSM 數據月度對帳 SOP：",
        "  - 每月 5 號前勾稽財務實收金流與 CSM 新購續約名單，將誤差控制在 5% 以內",
        "2. 優化 PLG 產品 Onboarding 引導流程：",
        "  - 針對逾半數註冊未簽署用戶，導入 3 分鐘動畫或範本。對註冊 48 小時內未動作者自動 Email 觸發並安排電訪",
        "3. 擴大對點點簽受災戶之精勢行銷：",
        "  - 行銷業務聚焦旅行社、連鎖醫療、人資派遣等，定向投放「不限件數吃到飽」廣告，加速轉單轉化"
    ]
    for item in pre_body:
        p = tfr_e.add_paragraph()
        p.text = item
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(12.5)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(6)

    # 4. 存檔 (考慮鎖定重試)
    try:
        prs.save(pptx_filepath)
        print(f"[SUCCESS] PPTX presentation generated at: {pptx_filepath}")
        print(f"File size: {os.path.getsize(pptx_filepath)} bytes")
        return True
    except PermissionError:
        print(f"[WARNING] {pptx_filename} is locked. Saving to alternates...")
        saved = False
        v = 2
        while not saved and v < 100:
            alt_path = pptx_filepath.replace(".pptx", f"_v{v}.pptx")
            try:
                prs.save(alt_path)
                print(f"[SUCCESS] PPTX successfully saved to: {alt_path}")
                saved = True
            except PermissionError:
                v += 1
        return saved
    except Exception as e:
        print(f"[ERROR] Failed to save PPTX: {e}")
        return False

if __name__ == "__main__":
    create_report_pptx()
