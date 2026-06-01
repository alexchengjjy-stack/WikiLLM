# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# 1. 定義顏色 (對齊 BreezySign CIS)
C_PRIMARY = RGBColor(5, 120, 87)       # 翠綠 (Primary Emerald)
C_SECONDARY = RGBColor(2, 132, 199)   # 天藍 (Secondary Sky)
C_LBG = RGBColor(236, 253, 245)       # 極淺綠 (Light Background)
C_DARK = RGBColor(15, 23, 42)         # Slate 深藍灰 (Dark Text)
C_WHITE = RGBColor(255, 255, 255)
C_GRAY_LIGHT = RGBColor(241, 245, 249) # Table header background

def main():
    workspace_dir = r"c:\Users\alexc\OneDrive\文件\WikiLLM"
    outputs_dir = os.path.join(workspace_dir, "outputs")
    template_path = os.path.join(outputs_dir, "templates", "bzs-presentation-template.pptx")
    
    if not os.path.exists(template_path):
        print(f"[ERROR] Template PPTX not found at: {template_path}")
        return False

    # 2. 定義時間戳記與檔名 (對齊 PDF)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M")
    pptx_filename = f"bzs-esign-monitoring-snapshot-202606-{timestamp}-v1.pptx"
    pptx_filepath = os.path.join(outputs_dir, "bzs", pptx_filename)

    # 3. 載入簡報 (保留格式與屬性)
    prs = Presentation(template_path)
    
    # 4. 刪除原模板簡報中所有的 mock slides
    for i in range(len(prs.slides)-1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Helper: 設定背景顏色
    def set_slide_bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    # Helper: 繪製圓角卡片
    def add_card(slide, left, top, width, height, bg_color, border_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # Helper: 繪製實心矩形
    def add_rect(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # Helper: 新增標準內容頁標頭 (CIS 線條與 Logo)
    def add_slide_header(slide, title_text):
        # 頂部翠綠線條
        add_rect(slide, 0, 0, Inches(13.33), Inches(0.15), C_PRIMARY)
        # 天藍輔助線
        add_rect(slide, 0, Inches(0.15), Inches(13.33), Inches(0.04), C_SECONDARY)
        
        # 標題文字區
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.5), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY

        # 右上角置入官方 Logo
        logo_green_path = os.path.join(outputs_dir, "assets", "bzs-logo-green.png")
        if os.path.exists(logo_green_path):
            slide.shapes.add_picture(logo_green_path, Inches(10.5), Inches(0.38), width=Inches(2.0), height=Inches(0.4))
        return title_box

    # =========================================================
    # Slide 1: 封面頁 (Cover Slide)
    # =========================================================
    slide_cover = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_cover, C_PRIMARY)

    # 頂部天藍裝飾塊
    add_rect(slide_cover, Inches(0), Inches(0), Inches(13.33), Inches(0.2), C_SECONDARY)

    # 插入官方高保真反白 Logo
    logo_white_path = os.path.join(outputs_dir, "bzs-logo-white.png")
    if os.path.exists(logo_white_path):
        slide_cover.shapes.add_picture(logo_white_path, Inches(1.2), Inches(1.8), width=Inches(2.61), height=Inches(0.52))

    # 封面主標題
    title_box = slide_cover.shapes.add_textbox(Inches(1.2), Inches(2.9), Inches(11.0), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "BreezySign 好好簽"
    p1.font.name = 'Microsoft JhengHei'
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "2026 年 6 月電子簽章能量登錄競品情報普查快照"
    p2.font.name = 'Microsoft JhengHei'
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_before = Pt(12)

    p3 = tf.add_paragraph()
    p3.text = "依據《普查與情報快照實測規範》對點點簽、律果簽、全景及我方 Production 正式官網 100% 實地查找與情報對照分析"
    p3.font.name = 'Microsoft JhengHei'
    p3.font.size = Pt(11)
    p3.font.color.rgb = RGBColor(200, 230, 215)
    p3.space_before = Pt(8)

    p4 = tf.add_paragraph()
    p4.text = "蒙恬科技 ． 電子簽章業務與技術整合小組"
    p4.font.name = 'Microsoft JhengHei'
    p4.font.size = Pt(13)
    p4.font.color.rgb = C_WHITE
    p4.space_before = Pt(35)

    # =========================================================
    # Slide 2: 四大電子簽章最新普查矩陣
    # =========================================================
    slide_matrix = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_matrix, C_WHITE)
    add_slide_header(slide_matrix, "一、 四大電子簽章最新普查矩陣 (2026年6月版)")

    # 新增表格
    rows, cols = 5, 5
    left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.73), Inches(4.8)
    table_shape = slide_matrix.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # 設定列寬
    table.columns[0].width = Inches(1.6)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.7)
    table.columns[3].width = Inches(2.5)
    table.columns[4].width = Inches(2.43)

    # 表格資料
    table_data = [
        ["數據維度", "好好簽 (BreezySign) 我方", "點點簽 (DottedSign) 競品", "律果簽 (LegalSign) 競品", "FastSIGN (全景) 競品"],
        ["能量登錄證號", "113電簽0008 (蒙恬科技)\n有效至民國 116/08/13", "113電簽0003 (凱鈿行動)\n有效至民國 115/08/29", "113電簽0005 (律果科技)\n有效至民國 116/01/15", "113電簽0001 (全景軟體)\n有效至民國 115/07/21"],
        ["產品動態與系統整合", "證號聲明正式同步上線官網。\n持續優化鼎新 ERP 整合，調優開啟連結逾時限制。", "1. 2026-05 整合 Vital BizForm\n2. 導入 MCP 大模型流程語意互動支援。", "大力推廣 AI 法務助理「法樂多」，主打 30秒自動審約、草擬與風險預判。", "1. 2026-03 推出 IDExpert Cloud\n2. 完成零信任三階段驗證，推廣 PQC 遷移。"],
        ["主要弱點與痛點", "官網 FAQ 模組大模型抓取仍為空白，亟需優化 GEO 能見度。", "2026-04-21 強制廢止舊企業方案續約，Envelope 計費引發跳船。", "SaaS 平台大批量合約 Loading 過慢（旅行社客戶實測需 10 分鐘）。", "極少投放線上行銷廣告，大眾級 GEO/SEO 能見度極低。"],
        ["我方作戰方針", "定價頁部署 6-8 組 FAQ 模組；投放防禦型品牌關鍵字廣告。", "集火 4/21 舊合約失效與漲價痛點，主打我方「年約吃到飽」方案。", "主打好好簽極速簽署效能（3秒完簽）與 LINE/簡訊傳簽靈活性。", "鎖定中大型企業，主打我方完整 AATL 與國家級能量登錄許可。"]
    ]

    for row_idx, row in enumerate(table_data):
        for col_idx, text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Microsoft JhengHei'
            
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_PRIMARY if col_idx == 1 else (C_SECONDARY if col_idx > 1 else C_DARK)
                p.font.bold = True
                p.font.size = Pt(11)
                p.font.color.rgb = C_WHITE
                p.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE if row_idx % 2 != 0 else C_GRAY_LIGHT
                p.font.size = Pt(9.5)
                p.font.color.rgb = C_DARK
                if col_idx == 0:
                    p.alignment = PP_ALIGN.CENTER
                    p.font.bold = True
                else:
                    p.alignment = PP_ALIGN.LEFT
                    if col_idx == 1:
                        p.font.color.rgb = C_PRIMARY

    # =========================================================
    # Slide 3: 核心競品情報通道深度解析 (1/2)
    # =========================================================
    slide_channel1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_channel1, C_WHITE)
    add_slide_header(slide_channel1, "二、 競品情報通道深度解析 (官網定價、策略文章、廣告字)")

    # 三欄卡片
    width_card = Inches(3.6)
    height_card = Inches(4.8)
    top_card = Inches(1.6)

    # 1. 官網與定價動向
    add_card(slide_channel1, Inches(0.8), top_card, width_card, height_card, C_LBG, C_PRIMARY)
    box1 = slide_channel1.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(3.4), Inches(4.5))
    tf1 = box1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "💰 官網與定價動向"
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY
    
    bullets1 = [
        "點點簽官網新計費落地：自 2026-04-21 起，舊版企業方案終止續約。到期強制轉為 Envelope 件數計費，否則降為免費停簽。",
        "律果簽加碼 AI CLM：基本版 $180、標準版 $490、專業版 $980，加掛數發部能量登錄許可章，並增設 AI 助理與 CLM 加購。",
        "全景主打買斷制與金融：行銷側重 FastSIGN Cloud 與 Pro 買斷方案，少宣傳訂閱。"
    ]
    for b in bullets1:
        p = tf1.add_paragraph()
        p.text = b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(11)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 2. 內容策略文章
    add_card(slide_channel1, Inches(4.8), top_card, width_card, height_card, RGBColor(240, 248, 255), C_SECONDARY)
    box2 = slide_channel1.shapes.add_textbox(Inches(4.9), Inches(1.7), Inches(3.4), Inches(4.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "✍️ 內容策略文章方向"
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_SECONDARY
    
    bullets2 = [
        "點點簽鎖定 ESG 與 HR：高頻產出「富士軟片 ESG 無紙化」與「人司入職自動化」文章，全數埋設 FAQ Schema 搶占 GEO。",
        "律果簽主攻 AI 與合約：Blog 全力轉型「AI 智慧合約講堂」，打造 AI 線上審約及法規風險預判的專業法務形象。",
        "全景偏向資安與零信任：主攻「製造業零信任資安升級」與「PQC 設備防護」，電簽 Blog 近兩月無更新。"
    ]
    for b in bullets2:
        p = tf2.add_paragraph()
        p.text = b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(11)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 3. 廣告策略與預算
    add_card(slide_channel1, Inches(8.8), top_card, width_card, height_card, C_GRAY_LIGHT, C_DARK)
    box3 = slide_channel1.shapes.add_textbox(Inches(8.9), Inches(1.7), Inches(3.4), Inches(4.5))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "🎯 廣告投放與關鍵字"
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_DARK
    
    bullets3 = [
        "點點簽狂砸廣告：預估月預算 NT$8萬-12萬，強力投放「電子簽章」及 API 詞，甚至買我方品牌字「BreezySign 費用」攔截。",
        "律果簽集火長尾高意圖字：月預算約 NT$2.5萬-4萬，投放「合約管理系統」、「線上審約」等字組。",
        "全景與我方投放空白：全景僅月投 NT$5K 且無多媒體廣告，主要依靠關係網；我方主要投放通用詞，品牌防禦依舊是空白。"
    ]
    for b in bullets3:
        p = tf3.add_paragraph()
        p.text = b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(11)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # =========================================================
    # Slide 4: 核心競品情報通道深度解析 (2/2)
    # =========================================================
    slide_channel2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_channel2, C_WHITE)
    add_slide_header(slide_channel2, "三、 競品情報通道深度解析 (人才招募、認證合規、公關及技術 SEO)")

    # 左右雙卡片
    width_half = Inches(5.6)
    height_half = Inches(4.8)
    top_half = Inches(1.6)

    # 左卡：人才與認證
    add_card(slide_channel2, Inches(0.8), top_half, width_half, height_half, C_LBG, C_PRIMARY)
    box_l = slide_channel2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.4))
    tfl_2 = box_l.text_frame
    tfl_2.word_wrap = True
    p = tfl_2.paragraphs[0]
    p.text = "👥 人才戰略與合規認證"
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    bullets_l = [
        "點點簽朝 AI 科技轉型： Yourator 釋出 AI Solution Consultant 與 Engineer 職缺，脫離單一工具定位。",
        "律果簽卡位法律 AI 專家：招聘「法律 AI 解決方案工程師」，要求具備 LLM 開發與法律條文 Fine-tuning 實務力。",
        "認證合規對齊：點點簽 (113電簽0003，有效至115/08/29)、律果簽 (113電簽0005，有效至116/01/15)；我方蒙恬科技 (113電簽0008，有效至116/08/13)，有效期為四大廠最長！"
    ]
    for b in bullets_l:
        p = tfl_2.add_paragraph()
        p.text = b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(12)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # 右卡：公關與技術 SEO
    add_card(slide_channel2, Inches(6.9), top_half, width_half, height_half, RGBColor(240, 248, 255), C_SECONDARY)
    box_r = slide_channel2.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.4))
    tfr_2 = box_r.text_frame
    tfr_2.word_wrap = True
    p = tfr_2.paragraphs[0]
    p.text = "🌐 公關新聞與技術 SEO/GEO"
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_SECONDARY

    bullets_r = [
        "點點簽公關權重高：發布「全球突破 2000 萬用戶」與數發部登錄公關稿，但因載入 HubSpot 行銷代碼導致 CWV INP 表現差。4/21 漲價亦導致大模型中出現客訴與負面引用漂移。",
        "律果簽公關盟軍多：發布與精誠資訊、數位時代等結盟公關，但在技術 SEO 方面因採用 SPA 架構且缺乏 SSR，易被誤判為薄內容，GEO 能見度差。",
        "全景零信任錨點強：2026-03 通過政府零信任三階段驗證。我方好好簽 Production 已正式上線能量登錄宣告，GEO 能見度達 40%，但 Organization Schema 仍可再指定精確。"
    ]
    for b in bullets_r:
        p = tfr_2.add_paragraph()
        p.text = b
        p.font.name = 'Microsoft JhengHei'
        p.font.size = Pt(12)
        p.font.color.rgb = C_DARK
        p.space_before = Pt(8)

    # =========================================================
    # Slide 5: 業務前線 Battle Cards 反駁小卡
    # =========================================================
    slide_battle = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_battle, C_WHITE)
    add_slide_header(slide_battle, "四、 業務前線即戰力：最新反駁小卡 (Battle Cards)")

    # 左右卡片
    # 左卡 (對決點點簽)
    add_card(slide_battle, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0), C_LBG, C_PRIMARY)
    left_b_box = slide_battle.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.6))
    tfl_b = left_b_box.text_frame
    tfl_b.word_wrap = True
    p = tfl_b.paragraphs[0]
    p.text = "🆚 當客戶提及「點點簽 (DottedSign)」時："
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    p_l1 = tfl_b.add_paragraph()
    p_l1.text = "【核心反擊點：強制漲價與件數計費】"
    p_l1.font.name = 'Microsoft JhengHei'
    p_l1.font.size = Pt(13)
    p_l1.font.bold = True
    p_l1.font.color.rgb = C_PRIMARY
    p_l1.space_before = Pt(8)

    p_l2 = tfl_b.add_paragraph()
    p_l2.text = "「點點簽自 2026 年 4 月 21 日起正式終止了舊版企業方案的續約，合約到期後強制轉換為以『任務件數』計費的新方案，這會造成您的簽署費用暴增。好好簽擁有完全對等的數位發展部『113電簽0008』核可與 AATL 國際憑證，且我們採用吃到飽方案與透明價格。我們支持在地化的 LINE 傳簽與聲明錄影防賴，背靠蒙恬科技 (30年大廠)，是法規與成本考量下的最穩健選擇。」"
    p_l2.font.name = 'Microsoft JhengHei'
    p_l2.font.size = Pt(12)
    p_l2.font.color.rgb = C_DARK
    p_l2.space_before = Pt(8)

    # 右卡 (對決律果簽)
    add_card(slide_battle, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.0), RGBColor(255, 253, 245), RGBColor(217, 119, 6))
    right_b_box = slide_battle.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.0), Inches(4.6))
    tfr_b = right_b_box.text_frame
    tfr_b.word_wrap = True
    p = tfr_b.paragraphs[0]
    p.text = "🆚 當客戶提及「律果簽 (LegalSign)」時："
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(185, 28, 28)

    p_r1 = tfr_b.add_paragraph()
    p_r1.text = "【核心反擊點：系統 Loading 性能與大批量簽署】"
    p_r1.font.name = 'Microsoft JhengHei'
    p_r1.font.size = Pt(13)
    p_r1.font.bold = True
    p_r1.font.color.rgb = RGBColor(185, 28, 28)
    p_r1.space_before = Pt(8)

    p_r2 = tfr_b.add_paragraph()
    p_r2.text = "「律果簽主打法律 AI 審約，但在電簽系統的底層性能上，大批量簽署常出現卡頓（旅行社大批量實測需 10 分鐘 Loading），極度延誤效率。好好簽核心引擎經過高頻壓力測試，穩定極速，且首創 LINE/簡訊傳簽與聲明錄影防賴。如果您需要的是一個極速流暢、具備法庭強證據力，且通過國家級能量登錄雙重核可的電子簽章平台，好好簽絕對是更實在、更好用的選擇。」"
    p_r2.font.name = 'Microsoft JhengHei'
    p_r2.font.size = Pt(12)
    p_r2.font.color.rgb = C_DARK
    p_r2.space_before = Pt(8)

    # =========================================================
    # Slide 6: 下一步行動方針
    # =========================================================
    slide_next = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_next, C_WHITE)
    add_slide_header(slide_next, "五、 好好簽下一步行銷與業務行動方針")

    # 行動方針
    next_box = slide_next.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5))
    tf_next = next_box.text_frame
    tf_next.word_wrap = True

    items = [
        ("1. [行銷/研發] 實作好好簽定價頁面 FAQ 模組", "比照點點簽於價格頁下方部署 6~8 組結構化問答 FAQ，方便大模型 (GEO/AEO) 進行爬取與定價意圖對照推薦。"),
        ("2. [行銷] 部署對手著陸頁面 (Competitor Landing Page)", "集火點點簽 4/21 舊合約強制停用與漲價痛點，主打好好簽「吃到飽方案」，精準搶攻流失之年簽中大戶。"),
        ("3. [廣告投放] 啟動品牌關鍵字廣告防禦與長尾合規詞", "立即投放防禦型字組，避免搜尋我方品牌「BreezySign」時被點點簽的搶客攔截廣告引流。")
    ]

    for title, desc in items:
        p_t = tf_next.add_paragraph() if tf_next.paragraphs[0].text else tf_next.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Microsoft JhengHei'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = C_PRIMARY
        p_t.space_before = Pt(15) if tf_next.paragraphs[0].text else Pt(0)

        p_d = tf_next.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Microsoft JhengHei'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = C_DARK
        p_d.space_before = Pt(5)

    # 5. 存檔 (考慮鎖定重試)
    try:
        prs.save(pptx_filepath)
        print(f"[SUCCESS] PPTX presentation generated at: {pptx_filepath}")
        print(f"File size: {os.path.getsize(pptx_filepath)} bytes")
        return True
    except PermissionError:
        print(f"[WARNING] {pptx_filename} is locked. Saving to alternates...")
        saved = False
        v = 2
        while not saved and v < 100:
            alt_path = pptx_filepath.replace(".pptx", f"_v{v}.pptx")
            try:
                prs.save(alt_path)
                print(f"[SUCCESS] PPTX successfully saved to: {alt_path}")
                saved = True
            except PermissionError:
                v += 1
        return saved
    except Exception as e:
        print(f"[ERROR] Failed to save PPTX: {e}")
        return False

if __name__ == "__main__":
    main()
