import os
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. 初始化與CIS配色
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9
prs.slide_height = Inches(7.5)

C_PRIMARY = RGBColor(5, 120, 87)       # 翠綠 (Primary Emerald)
C_SECONDARY = RGBColor(2, 132, 199)   # 天藍 (Secondary Sky)
C_LBG = RGBColor(236, 253, 245)       # 極淺綠 (Light Background)
C_DARK = RGBColor(15, 23, 42)         # Slate 深藍灰 (Dark Text)
C_WHITE = RGBColor(255, 255, 255)
C_GRAY_LIGHT = RGBColor(241, 245, 249) # 淺灰背景

# 2. 確定存檔路徑
outputs_dir = r"c:\Users\alexc\OneDrive\文件\WikiLLM\outputs"
timestamp = datetime.datetime.now().strftime("%Y%m%d-%H%M")
pptx_filename = f"bzs-2026h2-cross-department-plan-{timestamp}-v1.pptx"
pptx_filepath = os.path.join(outputs_dir, "bzs", pptx_filename)

os.makedirs(os.path.join(outputs_dir, "bzs"), exist_ok=True)

# 3. 輔助函數
def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_card(slide, left, top, width, height, bg_color, border_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

def add_slide_header(slide, title_text):
    # CIS 綠線
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(0.12), Inches(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_PRIMARY
    shape.line.fill.background()

    # 標題文字
    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.35), Inches(9.0), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    # 右上角綠色 Logo
    logo_green_path = os.path.join(outputs_dir, "assets", "bzs-logo-green.png")
    if os.path.exists(logo_green_path):
        slide.shapes.add_picture(logo_green_path, Inches(10.5), Inches(0.38), width=Inches(2.0), height=Inches(0.4))

# 4. 投影片製作
blank_layout = prs.slide_layouts[6]

# =========================================================
# Slide 1: 封面頁 (翠綠底 + 反白 Logo)
# =========================================================
slide_cover = prs.slides.add_slide(blank_layout)
set_slide_bg(slide_cover, C_PRIMARY)

logo_white_path = os.path.join(outputs_dir, "assets", "bzs-logo-white.png")
if os.path.exists(logo_white_path):
    slide_cover.shapes.add_picture(logo_white_path, Inches(1.2), Inches(1.8), width=Inches(2.61), height=Inches(0.52))

title_box = slide_cover.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(11.0), Inches(2.2))
tf = title_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "BreezySign 好好簽 ． 2026H2 跨部門執行計畫"
p1.font.name = 'Microsoft JhengHei'
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = C_WHITE

p2 = tf.add_paragraph()
p2.text = "PLG + SLG 雙軌成長戰術與財務、產品、技術合規行動方案"
p2.font.name = 'Microsoft JhengHei'
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(200, 230, 215)
p2.space_before = Pt(12)

p3 = tf.add_paragraph()
p3.text = "蒙恬科技 (PenPower) 好好簽跨部門成果驗證小組 ． 2026H2 (討論稿)"
p3.font.name = 'Microsoft JhengHei'
p3.font.size = Pt(12)
p3.font.italic = True
p3.font.color.rgb = C_WHITE
p3.space_before = Pt(40)

# =========================================================
# Slide 2: 核心戰略定位：PLG + SLG 雙軌引擎
# =========================================================
slide_strategy = prs.slides.add_slide(blank_layout)
set_slide_bg(slide_strategy, C_WHITE)
add_slide_header(slide_strategy, "一、 核心戰略定位：PLG 與 SLG 雙軌並進")

# 左卡 - PLG 戰術
add_card(slide_strategy, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0), C_LBG, C_PRIMARY)
left_box = slide_strategy.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.6))
tfl = left_box.text_frame
tfl.word_wrap = True

p = tfl.paragraphs[0]
p.text = "⚡ PLG (產品驅動成長) 戰術"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = C_PRIMARY

bullets_plg = [
    "引導自助轉化：鎖定中小企業與小微客群，著重低摩擦的自助式註冊與付費升級。",
    "克服冷啟動 (Cold Start)：透過優化 Onboarding 引導，目標將不活躍註冊率由目前的 50%~60% 降至 35% 以下。",
    "產品規格化限制：落實個人專業方案 (NT$3,000/年) 每月發送 100 份上限防線，降低變動憑證成本。",
    "超額加購憑證包：提供 100 份後超額加購 (NT$30/份)，享有高達 94.4% 的毛利率，將大戶化為利潤引擎。"
]
for b in bullets_plg:
    p = tfl.add_paragraph()
    p.text = b
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(12)
    p.font.color.rgb = C_DARK
    p.space_before = Pt(8)

# 右卡 - SLG 戰術
add_card(slide_strategy, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.0), RGBColor(240, 248, 255), C_SECONDARY)
right_box = slide_strategy.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.0), Inches(4.6))
tfr = right_box.text_frame
tfr.word_wrap = True

p = tfr.paragraphs[0]
p.text = "🤝 SLG (銷售驅動成長) 戰術"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = C_SECONDARY

bullets_slg = [
    "包抄競品大戶：集火每年合約簽署 5,000~20,000 份以上的大型客戶，對抗點點簽漲價與以件計費政策。",
    "專案無縫轉移特惠：針對競品流失之年簽大戶，主動提供首年 8 折或前 3 個月免費體驗，加速簽約決策。",
    "SI 通路對接：深化與百加 (101 BPM)、統一數網等 SI 通路的轉介合作，擴大業務觸角。",
    "鼎新電腦生態系整合：將電子簽章 API 深度整合至鼎新諸葛 AI 平台，鎖定製造與零售業核心流程。"
]
for b in bullets_slg:
    p = tfr.add_paragraph()
    p.text = b
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(12)
    p.font.color.rgb = C_DARK
    p.space_before = Pt(8)

# =========================================================
# Slide 3: 業務銷售與客戶成功 (Sales & CSM) 執行計畫
# =========================================================
slide_sales = prs.slides.add_slide(blank_layout)
set_slide_bg(slide_sales, C_WHITE)
add_slide_header(slide_sales, "二、 業務銷售與客戶成功 (Sales & CSM) 計畫")

# 核心指標卡
add_card(slide_sales, Inches(0.8), Inches(1.5), Inches(3.6), Inches(5.0), C_LBG, C_PRIMARY)
sales_metrics = slide_sales.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.2), Inches(4.6))
tf_sm = sales_metrics.text_frame
tf_sm.word_wrap = True

p = tf_sm.paragraphs[0]
p.text = "🎯 H2 業務核心指標"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = C_PRIMARY

metrics_body = [
    "● 付費公司數衝刺 1,750 家：",
    "  (5月實績為 397 家，需淨增 1,353 家，成長率 340.8%)",
    "● 舊客年約留存率維持 95% 以上：",
    "  (1月自動續約收入佔營收 42.4%，舊客維護乃穩定雪球)",
    "● 點點簽到期大戶銷售轉化率達 75% 以上",
    "● 台中浸信會等大戶簽訂「多年期預付款認列合約」"
]
for m in metrics_body:
    p = tf_sm.add_paragraph()
    p.text = m
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(11.5)
    p.font.color.rgb = C_DARK
    p.space_before = Pt(6)

# 執行策略
sales_strategies = slide_sales.shapes.add_textbox(Inches(4.8), Inches(1.5), Inches(7.7), Inches(5.0))
tf_ss = sales_strategies.text_frame
tf_ss.word_wrap = True

p = tf_ss.paragraphs[0]
p.text = "📋 業務核心執行策略"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = C_PRIMARY

strategies_body = [
    "1. 點點簽大戶無縫轉移特惠：提供首年 8 折或 3 個月體驗期，降低移轉決策阻力。",
    "2. 百加/鼎新通路商機對接：緊密跟進 ISV 通路，協同 SI 顧問費與分潤月結，理順合作機制。",
    "3. CSM 主動客戶防禦：對已開通用戶實施 Cold Start 二次啟用跟進，降低新註冊不活躍率。",
    "4. 大戶多年期長約收單：以 2~3 年長約提供固定費率優惠，鎖定客戶年繳金流，行政上安全收受預付款，並按會計年度逐年遞延認列為年租 MRR 營收。"
]
for s in strategies_body:
    p = tf_ss.add_paragraph()
    p.text = s
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(13)
    p.font.color.rgb = C_DARK
    p.space_before = Pt(12)

# =========================================================
# Slide 4: 行銷推廣 (Marketing) 執行計畫
# =========================================================
slide_mkt = prs.slides.add_slide(blank_layout)
set_slide_bg(slide_mkt, C_WHITE)
add_slide_header(slide_mkt, "三、 行銷推廣 (Marketing) H2 策略與成效")

# 左側 - 成效與漏斗
add_card(slide_mkt, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.0), RGBColor(240, 248, 255), C_SECONDARY)
mkt_left = slide_mkt.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(5.0), Inches(4.6))
tf_ml = mkt_left.text_frame
tf_ml.word_wrap = True

p = tf_ml.paragraphs[0]
p.text = "📈 實績指標與投資報酬"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = C_SECONDARY

mkt_metrics = [
    "● 超高投資報酬：窄口徑 LTV:CAC 達 67 倍，擴張空間極佳。",
    "● 廣告投放效益：2026 上半年品牌詞防禦 CPC 成本約 $20-$22 USD，轉換高，回報顯著。",
    "● 本土公信力槓桿：將數發部「113電簽0008」白名單背書轉為 GEO 引用之首要權威信號。",
    "● AI GEO 能見度突破：正式同步首頁與 Footer 聲明至正式站後，好好簽被列為 AI 搜尋推薦首選，能見度由 2.5 狂飆至 7.5+"
]
for m in mkt_metrics:
    p = tf_ml.add_paragraph()
    p.text = m
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(12)
    p.font.color.rgb = C_DARK
    p.space_before = Pt(8)

# 右側 - 執行策略
mkt_right = slide_mkt.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.0))
tf_mr = mkt_right.text_frame
tf_mr.word_wrap = True

p = tf_mr.paragraphs[0]
p.text = "📣 行銷核心執行戰術"
p.font.name = 'Microsoft JhengHei'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = C_SECONDARY

mkt_strategies = [
    "1. 對手著陸頁面 (Competitor Landing Page)：集火點點簽「件數計量與強制漲價」痛點，精準搶佔搜尋流失流量。",
    "2. 啟動品牌詞廣告防禦：立即針對「BreezySign」和「好好簽」投放防禦型廣告，防止流量被競品惡意攔截。",
    "3. 本地化合規行銷：主打中華電信 AATL 憑證、聲明錄影防賴等特色，突出本土資安防禦與台灣法規合規優勢。",
    "4. 公部門與 SI 聯合推案：製作公部門專用 EEAT 開發與 ISV 合作推廣文案，拓展 B2B 藍海。"
]
for s in mkt_strategies:
    p = tf_mr.add_paragraph()
    p.text = s
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(13)
    p.font.color.rgb = C_DARK
    p.space_before = Pt(10)

# =========================================================
# Slide 5: 產品經理與專案管理 (Product & PM) 規劃
# =========================================================
slide_prod = prs.slides.add_slide(blank_layout)
set_slide_bg(slide_prod, C_WHITE)
add_slide_header(slide_prod, "四、 產品經理與專案管理 (Product & PM) 規劃")

# 規格化三防線
prod_box = slide_prod.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.73), Inches(5.0))
tf_p = prod_box.text_frame
tf_p.word_wrap = True

products_body = [
    ("🛡️ 實施「大戶定價安全防線」與憑證加購包", "● 限制專業方案：當月發送合約額度硬性限制 100 份 / 月 (保護變動憑證成本)。\n● 超額加購包：達限制後阻斷，每份收取 NT$30 (最少購買5份)。此機制在常規混合場景下可創造高達 94.4% 的毛利率。"),
    ("🔑 重構「Unify 範本共用權限」角色分層", "● 權限分層設定：改變過去全開或全關模式，主帳號可將「範本管理與共享」授權給指定管理員。\n● 保留個人範本：子帳號既能共享企業範本，仍能自建私有的個人範本夾，完全滿足太平洋旅行社等多帳號情境。"),
    ("⚡ 優化「LINE/簡訊傳簽」極簡流程", "● 移除冗餘流程：免除傳簽時需「填入發起人手機號碼」之規定。\n● 一鍵複製連結：優化為一鍵生成並複製 LINE/簡訊簽署連結，在行動端減少輸入摩擦，大力吸引金融貸款及旅行社客群。")
]

for title, desc in products_body:
    p_t = tf_p.add_paragraph() if tf_p.paragraphs[0].text else tf_p.paragraphs[0]
    p_t.text = title
    p_t.font.name = 'Microsoft JhengHei'
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = C_PRIMARY
    p_t.space_before = Pt(12) if tf_p.paragraphs[0].text else Pt(0)

    p_d = tf_p.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Microsoft JhengHei'
    p_d.font.size = Pt(12.5)
    p_d.font.color.rgb = C_DARK
    p_d.space_before = Pt(4)

# =========================================================
# Slide 6: 技術工程 (RD / Engineering) 執行計畫
# =========================================================
slide_rd = prs.slides.add_slide(blank_layout)
set_slide_bg(slide_rd, C_WHITE)
add_slide_header(slide_rd, "五、 技術工程 (RD / Engineering) 執行計畫")

# RD 核心行動
rd_box = slide_rd.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.73), Inches(5.0))
tf_rd = rd_box.text_frame
tf_rd.word_wrap = True

rd_actions = [
    ("🌐 正式站首頁 DOM Heading 重構與 SEO/AEO 優化", "● DOM 大綱重構：首頁四大區塊主標題由 <h3> 修正為 <h2>，子描述順延為 <h3>，達成 HTML5 Outlining 完美巢狀，澄清爬蟲誤判，使正式站 Google PageSpeed 技術分數達 88 ~ 95 分。\n● 價格 Schema 同步：發布 Organization 與 Product JSON-LD 結構化資料，嵌入 NT$3,000 ~ 15,000 區間與評分格式，方便 AI 精準抓取。"),
    ("📂 公開表單「附件高品質打包下載」功能開發", "● 重構附件上傳：因應夢曦文化等大戶需求，除將身分證等附件合併於完簽 PDF 外，後台提供「一鍵單獨打包下載原始高品質證件」功能，以便司法存證對照與存查。"),
    ("⚙️ 自動化 API 測試沙盒 (Sandbox) 部署", "● 流程閉環開通：實現客戶申請 API 後，系統自動開通測試金鑰、自動發送 API 指南與 Postman Collection，大幅降低 RD 在測試支援的重複手動工時。")
]

for title, desc in rd_actions:
    p_t = tf_rd.add_paragraph() if tf_rd.paragraphs[0].text else tf_rd.paragraphs[0]
    p_t.text = title
    p_t.font.name = 'Microsoft JhengHei'
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = C_PRIMARY
    p_t.space_before = Pt(12) if tf_rd.paragraphs[0].text else Pt(0)

    p_d = tf_rd.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Microsoft JhengHei'
    p_d.font.size = Pt(12.5)
    p_d.font.color.rgb = C_DARK
    p_d.space_before = Pt(4)

# =========================================================
# Slide 7: 營運、行政與法務 (Ops & Legal) 合規修訂計畫
# =========================================================
slide_legal = prs.slides.add_slide(blank_layout)
set_slide_bg(slide_legal, C_WHITE)
add_slide_header(slide_legal, "六、 營運、行政與法務 (Ops & Legal) 合規修訂")

# ToS 與 隱私權修訂
legal_box = slide_legal.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.73), Inches(5.0))
tf_l = legal_box.text_frame
tf_l.word_wrap = True

legal_items = [
    ("⚖️ 用戶服務協議 (ToS) 升級要點 (對齊新版電子簽章法)", "● 數位簽章推定效力：增設「技術規格與法律效力宣告」專章，明定使用中華電信 AATL 或 TWCA 憑證方案即依法構成「數位簽章」，具有「推定本人親簽」效力，一般電子簽章則由用戶承擔否認舉證責任。\n● 默示合意機制：對齊新法第 5 條精神，刪除事前明示同意書面，更新為「若未明示反對且繼續簽署，即推定同意電子簽章」默示同意，降低傳簽摩擦力。\n● LINE 傳簽合意保障：於行動端點擊「確認送出簽章」即視為獲得合理反對機會且未反對，推定同意成立法律行為。"),
    ("🔒 隱私權政策 (Privacy Policy) 升級要點 (肖像個資與防線)", "● 生物特徵與錄影安全專章：針對「聲明錄影簽」敏感生物個資，宣告影像僅供司法存證特定目的利用，絕不用於廣告或演算法訓練，並承諾靜態 AES-256 加密與自動銷毀期限。\n● 憑證個資嵌入披露：明確告知使用者其姓名與遮蔽憑證代碼將永久且不可逆內嵌於完簽 PDF 的技術特性。\n● 表單附件個資責任分界：明確約定公開表單上傳附件之個資蒐集主體為「發起表單之企業客戶」，好好簽平台僅提供加密傳輸與安全儲存，不承擔實質利用與外洩之連帶法律責任。")
]

for title, desc in legal_items:
    p_t = tf_l.add_paragraph() if tf_l.paragraphs[0].text else tf_l.paragraphs[0]
    p_t.text = title
    p_t.font.name = 'Microsoft JhengHei'
    p_t.font.size = Pt(15)
    p_t.font.bold = True
    p_t.font.color.rgb = C_PRIMARY
    p_t.space_before = Pt(12) if tf_l.paragraphs[0].text else Pt(0)

    p_d = tf_l.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Microsoft JhengHei'
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = C_DARK
    p_d.space_before = Pt(4)

# =========================================================
# Slide 8: 2026H2 時程表與里程碑 (Milestones)
# =========================================================
slide_milestone = prs.slides.add_slide(blank_layout)
set_slide_bg(slide_milestone, C_WHITE)
add_slide_header(slide_milestone, "七、 2026H2 執行時程表與關鍵里程碑")

# 時程與里程碑
ms_box = slide_milestone.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.73), Inches(5.0))
tf_ms = ms_box.text_frame
tf_ms.word_wrap = True

milestones = [
    ("📅 第三季度 (Q3) - 產品定位與合規防線發布", "● 產品面：上線「大戶定價安全防線」每月 100 份軟性硬體限制，超額憑證加購包正式開售。\n● 行銷與業務：推廣「點點簽大戶無縫轉移特惠」，對接大瀚、得勝者、星鴻等大戶 Onboarding。\n● 法務與法學：正式發布新版合規 ToS 與隱私權政策個資安全特寫專章。"),
    ("📅 第四季度 (Q4) - API 自動化與生態通路全面引爆", "● 技術面：API 測試沙盒 (Sandbox) 標準化自動部署上線，公開表單「附件高品質打包下載」交付。\n● 通路面：與百加 BPM 完成系統轉介會計月結流程，深化統一數網經銷合作。\n● 行銷拓展：參加 9/18 新竹經濟部 AI Agent 應用供需媒合會上台分享，鼎新 ISV 諸葛平台對接與直播帶貨活動上線。"),
    ("📅 跨部門協同核心目標 - 衝刺付費盈利臨界點", "● 財務回報與盈利：淨增長 1,353 家付費客戶，推動好好簽由底數 397 家擴張至 1,750 家付費公司數的全面盈利爆發期。\n● 自動續訂營收雪球：舊客年繳留存率維持在 95% 以上，理順多年期預收款遞延 MRR 認列流程，打造穩固的利潤大後方。")
]

for title, desc in milestones:
    p_t = tf_ms.add_paragraph() if tf_ms.paragraphs[0].text else tf_ms.paragraphs[0]
    p_t.text = title
    p_t.font.name = 'Microsoft JhengHei'
    p_t.font.size = Pt(15.5)
    p_t.font.bold = True
    p_t.font.color.rgb = C_PRIMARY
    p_t.space_before = Pt(12) if tf_ms.paragraphs[0].text else Pt(0)

    p_d = tf_ms.add_paragraph()
    p_d.text = desc
    p_d.font.name = 'Microsoft JhengHei'
    p_d.font.size = Pt(12.5)
    p_d.font.color.rgb = C_DARK
    p_d.space_before = Pt(4)

# 5. 存檔
try:
    prs.save(pptx_filepath)
    print(f"[SUCCESS] PPTX presentation generated at: {pptx_filepath}")
    print(f"File size: {os.path.getsize(pptx_filepath)} bytes")
except PermissionError:
    print(f"[WARNING] {pptx_filename} is locked. Saving to alternates...")
    saved = False
    v = 2
    while not saved and v < 100:
        alt_path = pptx_filepath.replace(".pptx", f"_v{v}.pptx")
        try:
            prs.save(alt_path)
            print(f"[SUCCESS] PPTX successfully saved to: {alt_path}")
            saved = True
        except PermissionError:
            v += 1
except Exception as e:
    print(f"[ERROR] Failed to save PPTX: {e}")
