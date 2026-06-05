# -*- coding: utf-8 -*-
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def set_slide_background(slide, color):
    """設定投影片背景顏色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_card(slide, left, top, width, height, bg_color, border_color):
    """繪製高質感圓角卡片容器"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background() # 無邊框
    return shape

def add_slide_header(slide, title_text, category_text="COMPETITIVE INTELLIGENCE"):
    """加入標準投影片頁首標題 (大字體，霓虹藍)"""
    # 類別小字
    cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Outfit"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = RGBColor(188, 74, 255) # 霓虹紫
    
    # 頁首大標題 (放大至 32pt)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Microsoft JhengHei"
    p_title.font.size = Pt(32) # 大字體適合投影
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0, 229, 255) # 霓虹藍

def main():
    prs = Presentation()
    
    # 設定簡報為 16:9 寬螢幕標準尺寸
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6] # 乾淨空白佈局
    
    # 設計配色系統
    COLOR_BG = RGBColor(13, 17, 28)       # 科技暗黑底色 (Deep Slate Carbon)
    COLOR_CARD = RGBColor(23, 29, 43)     # 卡片深灰藍色 (Card Panel)
    COLOR_BORDER = RGBColor(40, 50, 75)   # 卡片邊框藍灰色
    COLOR_WHITE = RGBColor(240, 243, 248)  # 內文亮白灰色
    COLOR_MUTED = RGBColor(150, 160, 175)  # 內文暗灰色
    COLOR_CYAN = RGBColor(0, 229, 255)    # 霓虹藍 (標題)
    COLOR_PURPLE = RGBColor(188, 74, 255)  # 霓虹紫 (次要標題 / 亮點)
    COLOR_GREEN = RGBColor(0, 230, 118)   # 亮綠色 (優勢)
    COLOR_RED = RGBColor(255, 76, 76)     # 亮紅色 (痛點)
    
    # =========================================================================
    # SLIDE 1: 封面頁 (Cover Slide)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, COLOR_BG)
    
    # 背景大裝飾卡片
    create_card(slide1, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9), COLOR_CARD, COLOR_BORDER)
    
    # 封面大標題 (放大至 44pt)
    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "電子簽章能量登錄競品情報普查快照"
    p1.font.name = "Microsoft JhengHei"
    p1.font.size = Pt(44) # 封面大標題
    p1.font.bold = True
    p1.font.color.rgb = COLOR_CYAN
    
    # 封面副標題 (放大至 24pt)
    sub_box = slide1.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.8))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "BreezySign 2026 年 5 月首發基準線情報月報"
    p_sub.font.name = "Microsoft JhengHei"
    p_sub.font.size = Pt(24)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_PURPLE
    
    # 封面中繼資訊 (放大至 16pt)
    info_box = slide1.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.2))
    tf_info = info_box.text_frame
    
    p_info1 = tf_info.paragraphs[0]
    p_info1.text = "■ 普查執行時間：2026 年 5 月 14 日"
    p_info1.font.name = "Microsoft JhengHei"
    p_info1.font.size = Pt(16)
    p_info1.font.color.rgb = COLOR_WHITE
    
    p_info2 = tf_info.add_paragraph()
    p_info2.text = "■ 涵蓋監測競品：點點簽 (DottedSign) | 律果簽 (LegalSign) | 全景軟體 (Changingtec) | 好好簽 (BreezySign)"
    p_info2.font.name = "Microsoft JhengHei"
    p_info2.font.size = Pt(16)
    p_info2.font.color.rgb = COLOR_WHITE
    p_info2.space_before = Pt(8)

    # =========================================================================
    # SLIDE 2: 執行摘要與全局洞察 (Executive Summary)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, COLOR_BG)
    add_slide_header(slide2, "🌐 執行摘要與全局洞察 (Executive Summary)")
    
    # 左側：全局洞察主卡片 (寬 5.6)
    create_card(slide2, Inches(0.6), Inches(1.5), Inches(5.6), Inches(5.3), COLOR_CARD, COLOR_BORDER)
    left_box = slide2.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.0), Inches(4.7))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_lh = tf_left.paragraphs[0]
    p_lh.text = "產業全局大勢"
    p_lh.font.name = "Microsoft JhengHei"
    p_lh.font.size = Pt(20) # 卡片標題
    p_lh.font.bold = True
    p_lh.font.color.rgb = COLOR_PURPLE
    p_lh.space_after = Pt(14)
    
    p_lbody = tf_left.add_paragraph()
    p_lbody.text = "隨著數位發展部正式推動「電子簽章解決方案服務能量登錄」，國內電子簽章 SaaS 戰場已進入高度依賴「合規信任背書」與「深度生態系綁定」的成熟期。\n\n本期觀測顯示，單一電子簽署功能已難以滿足中大型企業客戶。對手正全面開闢「身分驗證行動化」與「自動化工作流串接」兩大主戰場，並積極招募法律與 AI 的雙棲複合型人才以構築技術護城河。"
    p_lbody.font.name = "Microsoft JhengHei"
    p_lbody.font.size = Pt(16) # 內文大字體
    p_lbody.font.color.rgb = COLOR_WHITE
    p_lbody.line_spacing = 1.3
    
    # 右側：三大產業集火趨勢 (寬 6.2)
    create_card(slide2, Inches(6.5), Inches(1.5), Inches(6.2), Inches(5.3), COLOR_CARD, COLOR_BORDER)
    right_box = slide2.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.7))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_rh = tf_right.paragraphs[0]
    p_rh.text = "三大產業集火趨勢"
    p_rh.font.name = "Microsoft JhengHei"
    p_rh.font.size = Pt(20) # 卡片標題
    p_rh.font.bold = True
    p_rh.font.color.rgb = COLOR_PURPLE
    p_rh.space_after = Pt(14)
    
    trends = [
        ("1. 身分驗證行動化 (Mobile Identity)", "點點簽大力整合「行動自然人憑證 (TWFidO)」，主打手機刷臉/指紋 10 秒內完成具不可否認性的數位簽章。"),
        ("2. 自動化流程串接 (Ecosystem Hooks)", "對手積極串接外部平台（如 SurveyCake 問卷串接），覆蓋合約前、中、後端完整生命週期，建立高度粘性。"),
        ("3. 法務科技雙棲人才戰 (LegalTech Talent)", "律果簽與全景軟體大舉開缺，鎖定結合「深度學習、影像識別、自然語言與法律實務」的跨界研發團隊。")
    ]
    
    for title, desc in trends:
        p_t = tf_right.add_paragraph()
        p_t.text = title
        p_t.font.name = "Microsoft JhengHei"
        p_t.font.size = Pt(16) # 趨勢標題 16pt 粗體
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CYAN
        p_t.space_before = Pt(8)
        
        p_d = tf_right.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Microsoft JhengHei"
        p_d.font.size = Pt(14) # 趨勢內容 14pt
        p_d.font.color.rgb = COLOR_WHITE
        p_d.space_after = Pt(6)

    # =========================================================================
    # SLIDE 3: 競品觀測基礎數據基準線 (PowerPoint Native Table)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, COLOR_BG)
    add_slide_header(slide3, "📊 競品觀測基礎數據基準線 (Market Baseline)")
    
    # 建立原生 PowerPoint 表格 (放大以利投影，行高與字體適度調整)
    rows = 8
    cols = 5
    left = Inches(0.5)
    top = Inches(1.4)
    width = Inches(12.333)
    height = Inches(5.4)
    
    table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # 設定各欄寬度
    table.columns[0].width = Inches(1.8)  # 維度
    table.columns[1].width = Inches(2.6)  # 點點簽
    table.columns[2].width = Inches(2.6)  # 律果簽
    table.columns[3].width = Inches(2.6)  # 好好簽
    table.columns[4].width = Inches(2.733) # FastSIGN (全景)
    
    # 表格資料內容 (新增月度 SEM 廣告預算數據)
    table_data = [
        ["數據維度", "點點簽 (DottedSign)", "律果簽 (LegalSign)", "好好簽 (BreezySign) 我方", "FastSIGN (全景)"],
        ["母公司 / 主體", "凱鈿行動科技 (Kdan)", "律果科技 (Legaltech)", "蒙恬科技 (PenPower)", "全景軟體 (Changingtec)"],
        ["員工人數規模", "全球約 200 - 230 人", "約 12 - 20 人 (純法務新創)", "總部約 90 人 (電簽獨立業務)", "約 170 人 (興櫃資安廠)"],
        ["預估月度廣告預算", "約 NT$ 80k - 120k / 月\n(強勢全通路買量)", "約 NT$ 25k - 40k / 月\n(精實高意圖長尾字)", "約 NT$ 15k - 25k / 月\n(偏通用免費流量字)", "約 NT$ 5k - 15k / 月\n(極低，主打線下招標)"],
        ["方案與價格", "• Pro: USD 15/月\n• Business: 計次任務包\n• Enterprise: 客製/API", "• 基本: NT$ 180/月\n• 標準: NT$ 490/月\n• 專業: NT$ 980/月\n• 企業: API客製", "• 免費: 3份/月\n• 專業: NT$ 300/月\n• 企業: NT$ 1,500/月\n• 客製: API 整合", "• 輕量: NT$ 25,000/年起\n• 專業: 地端買斷制\n• 政府/金融深度客製"],
        ["官方 Blog 篇數", "約 60 - 100 篇\n(高頻發布，含SEO長尾)", "約 30 - 50 篇\n(合約講堂、合約智庫)", "約 20 - 40 篇\n(電簽法、ESG、手寫)", "約 15 - 30 篇\n(資安公告、憑證技術)"],
        ["公開案例篇數", "約 9 - 15 篇\n( SurveyCake / YONEX )", "約 5 - 10 篇\n(精誠經銷、宏碁合約)", "官網無正式公開案例 (0篇)\n(內部日報已累積 10+ 家)", "約 15 - 25 篇\n(銀行、34+政府機關)"],
        ["核心優勢亮點", "國際品牌、API門檻低、大型任務包划算", "內建 CLM、合約全生命週期管理、法顧加持", "聲明錄影防賴強效證據、手寫板完美整合、LINE 傳簽", "地端買斷授權、純血資安大廠、AD系統深度整合"]
    ]
    
    # 填充表格內容並設定精美樣式
    for r_idx in range(rows):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            cell.text = table_data[r_idx][c_idx]
            
            # 設定儲存格背景色 (表頭 vs 資料行交錯底色)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(40, 50, 75) # 深藍灰表頭
            elif r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(23, 29, 43) # 交錯色卡片深
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(16, 21, 33) # 交錯色深黑
                
            # 字型設定
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.font.name = "Microsoft JhengHei"
                
                if r_idx == 0:
                    paragraph.font.size = Pt(12) # 表頭 12pt (稍微調小防止溢出)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLOR_CYAN
                else:
                    paragraph.font.size = Pt(10) # 表格內文 10pt，確保 8 行完整呈現不溢出
                    paragraph.font.color.rgb = COLOR_WHITE
                    # 亮顯我方或特殊點
                    if c_idx == 3: # 我方亮點
                        if r_idx == 6: # 0篇警告
                            paragraph.font.color.rgb = COLOR_RED
                            paragraph.font.bold = True
                        elif r_idx == 7: # 優勢
                            paragraph.font.color.rgb = COLOR_GREEN
                            paragraph.font.bold = True

    # =========================================================================
    # SLIDE 4: 競品解析：點點簽 vs. 律果簽 (DottedSign vs. LegalSign)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, COLOR_BG)
    add_slide_header(slide4, "⚔️ 競品解析：點點簽 vs. 律果簽")
    
    # 左側：點點簽 (DottedSign) — 領跑者的全通路防禦 (寬 5.8)
    create_card(slide4, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), COLOR_CARD, COLOR_BORDER)
    left_box = slide4.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.2), Inches(4.7))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "点点签 (DottedSign) — 領跑者防禦"
    p_lh.font.name = "Microsoft JhengHei"
    p_lh.font.size = Pt(20) # 標題
    p_lh.font.bold = True
    p_lh.font.color.rgb = COLOR_PURPLE
    p_lh.space_after = Pt(12)
    
    points_l = [
        ("🌐 完美的單頁 FAQ 部署 (GEO 友善)", "定價頁底部完整展現問答對 (無隱藏跳轉)，這是生成式 AI (GEO) 抓取標準答案的完美範本，AI 推薦引用率高達 85%。"),
        ("📝 高頻率內容行銷與 ESG 集火", "Blog 週週發文。近期主打 SurveyCake 問卷整合以覆蓋 HR 入職流程，以及 YONEX 導入電簽「降低 90% 行政工時」的 ESG 量化實證。"),
        ("👥 人才戰略轉向「AI 顧問化」", "Yourator 釋出 AI Solution Consultant 等高階職缺，產品正從「單一簽署工具」升級為「AI 驅動的文件解決方案」。")
    ]
    
    for title, desc in points_l:
        p_t = tf_l.add_paragraph()
        p_t.text = title
        p_t.font.name = "Microsoft JhengHei"
        p_t.font.size = Pt(15) # 項目標題大字體
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CYAN
        p_t.space_before = Pt(6)
        
        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Microsoft JhengHei"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = COLOR_WHITE
        p_d.space_after = Pt(6)
        
    # 右側：律果簽 (LegalSign) — 垂直深耕的法遵護城河 (寬 5.8)
    create_card(slide4, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3), COLOR_CARD, COLOR_BORDER)
    right_box = slide4.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.2), Inches(4.7))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "律果签 (LegalSign) — 法務技術壁壘"
    p_rh.font.name = "Microsoft JhengHei"
    p_rh.font.size = Pt(20) # 標題
    p_rh.font.bold = True
    p_rh.font.color.rgb = COLOR_PURPLE
    p_rh.space_after = Pt(12)
    
    points_r = [
        ("🌐 價格公開透明但存在「GEO 萃取硬傷」", "定價 NT$ 180 / 490 / 980 明碼標價；但常見問題僅有標題清單，需點擊跳轉才能閱讀，增加 AI 爬蟲抽取資訊時的斷鏈風險。"),
        ("⚖️ 強大的執業律師與合約智庫背書", "主打「合約生命週期管理 (CLM)」，資源高度集中於合約範本與律師起草服務，在回答合規與法律保障等 AI 提問時權重極高。"),
        ("👥 招募法律 AI 工程師築城", "Yourator 招募「法律 AI 解決方案工程師」，開發 LLM 法律助手與自動審約演算法，深耕垂直法務科技領域。")
    ]
    
    for title, desc in points_r:
        p_t = tf_r.add_paragraph()
        p_t.text = title
        p_t.font.name = "Microsoft JhengHei"
        p_t.font.size = Pt(15) # 項目標題大字體
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CYAN
        p_t.space_before = Pt(6)
        
        p_d = tf_r.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Microsoft JhengHei"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = COLOR_WHITE
        p_d.space_after = Pt(6)

    # =========================================================================
    # SLIDE 5: 競品解析：全景軟體 vs. 好好簽 (Changingtec vs. BreezySign)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, COLOR_BG)
    add_slide_header(slide5, "🛡️ 競品解析：全景軟體 vs. 好好簽我方現狀")
    
    # 左側：全景軟體 (FastSIGN) — 基礎架構巨頭的軟硬兼施 (寬 5.8)
    create_card(slide5, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), COLOR_CARD, COLOR_BORDER)
    left_box = slide5.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.2), Inches(4.7))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "全景軟體 (FastSIGN) — 資安與地端優勢"
    p_lh.font.name = "Microsoft JhengHei"
    p_lh.font.size = Pt(20)
    p_lh.font.bold = True
    p_lh.font.color.rgb = COLOR_PURPLE
    p_lh.space_after = Pt(12)
    
    points_l = [
        ("💼 獨特的「地端永久授權買斷」方案", "除了 FastSIGN 雲端 SaaS 方案，地端買斷與客製是其主戰場。深耕政府、大型金融單位的 PKI 憑證與零信任整合。"),
        ("📢 強大的大型標案與公信力背書", "常態性獲取公股行庫與縣市政府底層憑證系統標案，網頁案例高達 15-25 篇，在公部門端具有極高權威聲量。"),
        ("👥 影像識別與高階零信任佈局", "104 大舉招募「深度學習開發」、「影像軟體」與「資安/網安」工程師，著手研發基於電腦視覺的證件防偽比對與 OCR 審約。")
    ]
    
    for title, desc in points_l:
        p_t = tf_l.add_paragraph()
        p_t.text = title
        p_t.font.name = "Microsoft JhengHei"
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CYAN
        p_t.space_before = Pt(6)
        
        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Microsoft JhengHei"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = COLOR_WHITE
        p_d.space_after = Pt(6)
        
    # 右側：好好簽 (BreezySign) 我方現狀優勢 (寬 5.8)
    create_card(slide5, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3), COLOR_CARD, COLOR_BORDER)
    right_box = slide5.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.2), Inches(4.7))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "好好簽 (BreezySign) 我方現狀優勢"
    p_rh.font.name = "Microsoft JhengHei"
    p_rh.font.size = Pt(20)
    p_rh.font.bold = True
    p_rh.font.color.rgb = COLOR_GREEN # 綠色亮顯我方優勢
    p_rh.space_after = Pt(12)
    
    points_r = [
        ("📹 聲明錄影防賴強效證據 (全台唯一)", "支援「簽署過程錄音錄影」，為我方最具防非防賴力、加強法律保障的差異化招牌功能。"),
        ("🤝 蒙恬觸控面板整合硬體優勢", "與母公司蒙恬科技手寫硬體完美結合，專攻實體臨櫃面簽與臨櫃簽核，建立硬體護城河。"),
        ("💬 LINE 傳簽與極致在地化通知", "支援 LINE 直接發送合約與簽署通知，深度貼合台灣大眾與中小企業的社交與行動辦公習慣。"),
        ("🛡️ AATL 國際信賴清單憑證與印章儲存", "支援符合華人商業習慣的印章儲存管理，保障符合我國《電子簽章法》之推定效力與安全性。")
    ]
    
    for title, desc in points_r:
        p_t = tf_r.add_paragraph()
        p_t.text = title
        p_t.font.name = "Microsoft JhengHei"
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CYAN
        p_t.space_before = Pt(4)
        
        p_d = tf_r.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Microsoft JhengHei"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = COLOR_WHITE
        p_d.space_after = Pt(4)

    # =========================================================================
    # SLIDE 6: 我方痛點與落後指標 (BreezySign Gap & Threat)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, COLOR_BG)
    add_slide_header(slide6, "⚠️ 好好簽痛點診斷與戰略威脅")
    
    # 左側：我方痛點與落後指標 (寬 5.8)
    create_card(slide6, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), COLOR_CARD, COLOR_BORDER)
    left_box = slide6.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.2), Inches(4.7))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "好好簽我方痛點與落後指標"
    p_lh.font.name = "Microsoft JhengHei"
    p_lh.font.size = Pt(20)
    p_lh.font.bold = True
    p_lh.font.color.rgb = COLOR_RED # 紅色亮顯痛點
    p_lh.space_after = Pt(12)
    
    points_l = [
        ("🔴 核心公信力「官網宣佈缺席」", "**致命傷**：雖然我方已通過數發部「服務能量登錄」，但好好簽產品官網首頁、功能頁卻**隻字未提**，導致潛在客戶與大模型搜尋時遭遇「公信力靜默斷鏈」。"),
        ("🔴 GEO 萃取資訊黑洞 (FAQ 解答空白)", "價格常見問題模組「僅有問題、沒有文字解答內容」，大模型在提取對比定價與規則時面臨空白。"),
        ("🔴 官網案例數為 0，SEO 品牌發散", "官網缺乏任何 Case Studies 內容行銷；且「蒙恬好好簽」與「BreezySign」品牌詞混用，稀釋實體權重。")
    ]
    
    for title, desc in points_l:
        p_t = tf_l.add_paragraph()
        p_t.text = title
        p_t.font.name = "Microsoft JhengHei"
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CYAN
        p_t.space_before = Pt(6)
        
        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Microsoft JhengHei"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = COLOR_WHITE
        p_d.space_after = Pt(6)
        
    # 右側：產業競爭威脅與戰略反擊 (寬 5.8) (更新廣告費防禦策略)
    create_card(slide6, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.3), COLOR_CARD, COLOR_BORDER)
    right_box = slide6.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.2), Inches(4.7))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "戰略威脅與應對方針"
    p_rh.font.name = "Microsoft JhengHei"
    p_rh.font.size = Pt(20)
    p_rh.font.bold = True
    p_rh.font.color.rgb = COLOR_PURPLE
    p_rh.space_after = Pt(12)
    
    threats = [
        ("⚡ 凱鈿 (DottedSign) AI 顧問化威脅 [🔴 高]", "凱鈿透過 AI Solution 團隊將電簽工具包裝成大型解決方案，對抗通用型 AI 勢頭。\n→ **應對方案**：我方應強化「與鼎新 ERP 系統深度整合」的**即戰力論述**。"),
        ("⚡ 競品高預算廣告壓制與戰略突圍 [🔴 高]", "點點簽每月 10 萬廣告費壓制通用字，我方難以正面匹敵。\n→ **應對方案**：我方應集火攔截 `數發部能量登錄`、`符合電子簽章法` 等合規意圖長尾字，避開正面戰場。")
    ]
    
    for title, desc in threats:
        p_t = tf_r.add_paragraph()
        p_t.text = title
        p_t.font.name = "Microsoft JhengHei"
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CYAN
        p_t.space_before = Pt(8)
        
        p_d = tf_r.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Microsoft JhengHei"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = COLOR_WHITE
        p_d.space_after = Pt(8)

    # =========================================================================
    # SLIDE 7: 業務前線即戰力：銷售對話反駁指南 (Battle Cards)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7, COLOR_BG)
    add_slide_header(slide7, "🎯 業務前線 Battle Cards (銷售反駁話術)")
    
    # 兩側銷售話術卡片 (左右分欄)
    card_w = Inches(5.8)
    card_w = Inches(5.8)
    card_h = Inches(5.3)
    
    # 卡片 1: 對抗點點簽 (DottedSign)
    create_card(slide7, Inches(0.6), Inches(1.5), card_w, card_h, COLOR_CARD, COLOR_BORDER)
    box_1 = slide7.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.2), Inches(4.7))
    tf1 = box_1.text_frame
    tf1.word_wrap = True
    
    p_h1 = tf1.paragraphs[0]
    p_h1.text = "🆚 對抗點點簽 (DottedSign)"
    p_h1.font.name = "Microsoft JhengHei"
    p_h1.font.size = Pt(20)
    p_h1.font.bold = True
    p_h1.font.color.rgb = COLOR_CYAN
    p_h1.space_after = Pt(8)
    
    p_q1 = tf1.add_paragraph()
    p_q1.text = "● 客戶質疑：「點點簽市佔率高，而且串接了 SurveyCake 問卷，感覺很完整。」"
    p_q1.font.name = "Microsoft JhengHei"
    p_q1.font.size = Pt(14)
    p_q1.font.color.rgb = COLOR_MUTED
    p_q1.space_after = Pt(10)
    
    p_a1_title = tf1.add_paragraph()
    p_a1_title.text = "● 我方反駁話術框架 (與潛在客戶對話)："
    p_a1_title.font.name = "Microsoft JhengHei"
    p_a1_title.font.size = Pt(14)
    p_a1_title.font.bold = True
    p_a1_title.font.color.rgb = COLOR_GREEN
    
    p_a1 = tf1.add_paragraph()
    p_a1.text = "「點點簽做得很早，但其商務計費核心是基於『簽署任務次戰次數』，當您企業用印量或自動化發送量大增時，後續加購成本會直線上升。\n\n好好簽同樣通過數發部電子簽章服務能量登錄官方審查，具備同等最高等級數位憑證不可否認性；且我們背靠上櫃公司蒙恬科技 30 年的技術支援，能提供更具彈性且透明的授權計費結構，不會用次數綁架您的長遠數位化流程。」"
    p_a1.font.name = "Microsoft JhengHei"
    p_a1.font.size = Pt(14)
    p_a1.font.color.rgb = COLOR_WHITE
    p_a1.space_before = Pt(6)
    p_a1.line_spacing = 1.25

    # 卡片 2: 對抗律果簽 (LegalSign)
    create_card(slide7, Inches(6.9), Inches(1.5), card_w, card_h, COLOR_CARD, COLOR_BORDER)
    box_2 = slide7.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.2), Inches(4.7))
    tf2 = box_2.text_frame
    tf2.word_wrap = True
    
    p_h2 = tf2.paragraphs[0]
    p_h2.text = "🆚 對抗律果簽 (LegalSign)"
    p_h2.font.name = "Microsoft JhengHei"
    p_h2.font.size = Pt(20)
    p_h2.font.bold = True
    p_h2.font.color.rgb = COLOR_CYAN
    p_h2.space_after = Pt(8)
    
    p_q2 = tf2.add_paragraph()
    p_q2.text = "● 客戶質疑：「律果簽有律師團隊撰寫範本，價格方案也標得很清楚。」"
    p_q2.font.name = "Microsoft JhengHei"
    p_q2.font.size = Pt(14)
    p_q2.font.color.rgb = COLOR_MUTED
    p_q2.space_after = Pt(10)
    
    p_a2_title = tf2.add_paragraph()
    p_a2_title.text = "● 我方反駁話術框架 (與潛在客戶對話)："
    p_a2_title.font.name = "Microsoft JhengHei"
    p_a2_title.font.size = Pt(14)
    p_a2_title.font.bold = True
    p_a2_title.font.color.rgb = COLOR_GREEN
    
    p_a2 = tf2.add_paragraph()
    p_a2.text = "「律果簽的 CLM 與律師範本確實有其特色，但純軟新創在底層系統整合與長期營運穩定度上，往往不如具備軟硬體原生研發實力的老牌大廠。\n\n好好簽不僅提供符合電子簽章法推定效力的數位簽署，更整合了蒙恬科技深厚的影像與文字辨識底層技術。若您需要的是一套運行極速、直覺且具備國家級能量登錄資質的純粹電子簽核中樞，好好簽能給予您更輕量、無痛的導入體驗。」"
    p_a2.font.name = "Microsoft JhengHei"
    p_a2.font.size = Pt(14)
    p_a2.font.color.rgb = COLOR_WHITE
    p_a2.space_before = Pt(6)
    p_a2.line_spacing = 1.25

    # =========================================================================
    # SLIDE 8: 下一步行動與修復優先清單 (Next Steps)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8, COLOR_BG)
    add_slide_header(slide8, "🚀 下一步行動與修復優先清單 (Next Steps)")
    
    # 全寬大容器
    create_card(slide8, Inches(0.6), Inches(1.5), Inches(12.133), Inches(5.3), COLOR_CARD, COLOR_BORDER)
    next_box = slide8.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.533), Inches(4.7))
    tf_next = next_box.text_frame
    tf_next.word_wrap = True
    
    p_nh = tf_next.paragraphs[0]
    p_nh.text = "為了在下一期普查快照中實現 SEO/GEO 指標反超，研發與行銷團隊應啟動以下修復："
    p_nh.font.name = "Microsoft JhengHei"
    p_nh.font.size = Pt(18)
    p_nh.font.bold = True
    p_nh.font.color.rgb = COLOR_PURPLE
    p_nh.space_after = Pt(20)
    
    steps = [
        ("1. [行銷/前端] 立即實作定價頁面 FAQ 模組", 
         "參照點點簽格式，於好好簽定價頁面底部寫入 6~8 組有問有答的完整文字段落（而非展開列表），徹底解決 AI 大模型抓取空白、資訊黑洞的問題。"),
        ("2. [SEO/技術] 顯式掛載「能量登錄標章」與重構 Title Tag", 
         "在好好簽官網首頁及功能頁顯式宣布「通過數位發展部電子簽章解決方案服務能量登錄審查」；並重構 Title Tag 關鍵字，注入「電子簽名」、「電子合約管理」與「數發部能量登錄」。"),
        ("3. [廣告投放] 重啟合規長尾關鍵字廣告與攔截戰術", 
         "配置小額廣告預算，定向投放「數發部能量登錄」與「台灣電子簽章法合規」等合規長尾詞，攔截法人與大模型檢索意圖，突破競品通用高價字的壓制。")
    ]
    
    for title, desc in steps:
        p_t = tf_next.add_paragraph()
        p_t.text = title
        p_t.font.name = "Microsoft JhengHei"
        p_t.font.size = Pt(16) # 步驟標題 16pt 粗體
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CYAN
        p_t.space_before = Pt(10)
        
        p_d = tf_next.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Microsoft JhengHei"
        p_d.font.size = Pt(14) # 步驟說明 14pt
        p_d.font.color.rgb = COLOR_WHITE
        p_d.space_before = Pt(4)
        p_d.space_after = Pt(10)
        p_d.line_spacing = 1.2

    # 保存檔案至臨時目錄，以避免 Windows 中文編碼問題，後續將由 PowerShell 搬移
    output_path = r"C:\Users\alexc\AppData\Local\Temp\esign-monitoring-snapshot-202605.pptx"
    prs.save(output_path)
    print("Success: PPTX generated and saved to temporary folder.")

if __name__ == "__main__":
    main()
